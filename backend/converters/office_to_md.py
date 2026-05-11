from typing import List, Optional, Dict, Tuple
from collections import Counter
import os
import unicodedata
from .base_converter import BaseConverter
import re
import subprocess
import tempfile
from pathlib import Path
from datetime import datetime

# PDF处理 - PyMuPDF 优先
try:
    import fitz
    fitz_available = True
except ImportError:
    fitz_available = False

# PDF处理 - pypdf 作为 OCR 回退
try:
    import pypdf
    import pytesseract
    from pdf2image import convert_from_path
    pypdf_available = True
except ImportError:
    pypdf_available = False

# Office文件处理
try:
    import docx2txt
    docx_available = True
except ImportError:
    docx_available = False

try:
    import pandas as pd
    import tabulate
    pandas_available = True
except ImportError:
    pandas_available = False

try:
    from pptx import Presentation
    pptx_available = True
except ImportError:
    pptx_available = False

# HTML转Markdown
try:
    import html2text
    html2text_available = True
except ImportError:
    html2text_available = False

pdf_available = fitz_available or pypdf_available


class OfficeToMdConverter(BaseConverter):
    """
    Office 文档到 Markdown 转换器

    支持 PDF/DOCX/XLSX/PPTX/HTML -> Markdown 转换
    PDF 转换使用 PyMuPDF，根据字体大小和加粗样式自动识别标题层级
    """

    def __init__(self, output_dir: str, **kwargs):
        super().__init__(output_dir, **kwargs)
        self.poppler_path = kwargs.get('poppler_path')
        self.tesseract_cmd = kwargs.get('tesseract_cmd')
        self._check_dependencies()

    def _check_dependencies(self):
        missing_deps = []

        if not fitz_available:
            missing_deps.append("PyMuPDF (用于PDF智能转换，推荐安装)")
            self.logger.warning("PyMuPDF未安装，PDF转换将使用基础模式")

        if not pypdf_available:
            self.logger.info("pypdf/pytesseract/pdf2image未安装，OCR回退功能不可用")

        if not pdf_available:
            missing_deps.append("PyMuPDF 或 pypdf (至少需要一种PDF处理库)")
            self.logger.warning("PDF处理库均未安装，PDF转换功能将受限")

        if not docx_available:
            missing_deps.append("docx2txt (用于Word文档处理)")
            self.logger.warning("docx2txt库未安装，Word转换功能将受限")

        if not pandas_available:
            missing_deps.append("pandas, tabulate (用于Excel文件处理)")
            self.logger.warning("pandas/tabulate库未安装，Excel转换功能将受限")

        if not pptx_available:
            missing_deps.append("python-pptx (用于PowerPoint文件处理)")
            self.logger.warning("python-pptx库未安装，PowerPoint转换功能将受限")

        if not html2text_available:
            self.logger.warning("html2text库未安装，HTML转换功能将受限")

        try:
            import psutil
            self._psutil_available = True
        except ImportError:
            self._psutil_available = False
            missing_deps.append("psutil (用于批量并行处理和内存监控，推荐安装)")
            self.logger.info("psutil未安装，批量并行处理功能不可用，将使用串行处理")

        try:
            subprocess.run(["tesseract", "--version"],
                          stdout=subprocess.PIPE,
                          stderr=subprocess.PIPE,
                          check=True)
        except (subprocess.SubprocessError, FileNotFoundError):
            missing_deps.append("tesseract-ocr (用于图像OCR处理)")
            self.logger.warning("tesseract未安装，OCR功能将不可用")

        if missing_deps:
            self.logger.warning(f"以下依赖缺失，部分功能可能不可用: {', '.join(missing_deps)}")

    def convert(self, input_path: str) -> List[str]:
        supported_extensions = ['.pdf', '.docx', '.doc', '.xlsx', '.xls', '.xlsm',
                              '.pptx', '.ppt', '.html', '.htm']

        if not self._is_valid_input(input_path, supported_extensions):
            raise ValueError(f"无效的输入文件或目录: {input_path}")

        output_files = []

        if os.path.isfile(input_path):
            output_file = self._convert_single_file(input_path)
            if output_file:
                output_files.append(output_file)
        else:
            office_files = self._get_files_by_extension(input_path, supported_extensions)
            if not office_files:
                raise ValueError(f"目录中未找到支持的Office文件: {input_path}")

            for office_file in office_files:
                output_file = self._convert_single_file(office_file)
                if output_file:
                    output_files.append(output_file)

        return output_files

    def _convert_single_file(self, file_path: str) -> Optional[str]:
        file_path_obj = Path(file_path)
        file_type = self._get_file_type(file_path_obj)

        if not file_type:
            self.logger.warning(f"不支持的文件类型: {file_path}")
            return None

        try:
            self.logger.info(f"正在处理{file_type}文件: {file_path_obj.name}...")

            if file_type == 'pdf':
                md_text = self._extract_text_from_pdf(file_path_obj)
            elif file_type == 'word':
                text = self._extract_text_from_word(file_path_obj)
                md_text = self._convert_to_markdown(text) if text else ""
            elif file_type == 'excel':
                md_text = self._extract_text_from_excel(file_path_obj)
            elif file_type == 'powerpoint':
                md_text = self._extract_text_from_powerpoint(file_path_obj)
            elif file_type == 'html':
                md_text = self._extract_text_from_html(file_path_obj)
            else:
                self.logger.warning(f"未知文件类型: {file_type}")
                return None

            if md_text:
                md_text = self._optimize_markdown(md_text)
                output_file = self._save_markdown(md_text, file_path_obj)
                return output_file
            else:
                self.logger.warning(f"无法从 {file_path_obj.name} 提取文本")
                return None

        except Exception as e:
            self.logger.error(f"处理文件 {file_path} 失败: {str(e)}")
            return None

    def _get_file_type(self, file_path: Path) -> Optional[str]:
        suffix = file_path.suffix.lower()

        if suffix in ['.pdf']:
            return 'pdf'
        elif suffix in ['.docx', '.doc']:
            return 'word'
        elif suffix in ['.xlsx', '.xls', '.xlsm']:
            return 'excel'
        elif suffix in ['.pptx', '.ppt']:
            return 'powerpoint'
        elif suffix in ['.html', '.htm']:
            return 'html'
        else:
            return None

    # ─────────────────────────────────────────────
    # PDF 提取 - PyMuPDF 智能转换
    # ─────────────────────────────────────────────

    def _extract_text_from_pdf(self, pdf_path: Path) -> str:
        if not pdf_available:
            self.logger.error("PDF处理库未安装，无法处理PDF文件")
            return "## PDF内容提取失败\n\n需安装PDF处理库: pip install PyMuPDF"

        if fitz_available:
            return self._extract_pdf_with_fitz(pdf_path)
        else:
            return self._extract_pdf_with_pypdf(pdf_path)

    def _extract_pdf_with_fitz(self, pdf_path: Path) -> str:
        try:
            doc = fitz.open(pdf_path)
        except Exception as e:
            self.logger.error(f"无法打开PDF文件 {pdf_path}: {e}")
            return ""

        try:
            assets_dir_name = f"{pdf_path.stem}_assets"
            assets_dir = Path(self.output_dir) / assets_dir_name
            assets_dir_created = False

            all_page_spans = []
            page_heights = []
            page_count = len(doc)

            for page_idx in range(page_count):
                page = doc[page_idx]
                page_dict = page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE)
                page_spans = []
                page_heights.append(page.rect.height)

                for block in page_dict.get("blocks", []):
                    if block.get("type") != 0:
                        continue
                    block_bbox = block.get("bbox", (0, 0, 0, 0))
                    for line in block.get("lines", []):
                        line_text = ""
                        line_spans = []
                        for span in line.get("spans", []):
                            text = span.get("text", "").strip()
                            if not text:
                                continue
                            line_text += span.get("text", "")
                            line_spans.append({
                                "text": text,
                                "size": round(span.get("size", 12), 1),
                                "font": span.get("font", ""),
                                "flags": span.get("flags", 0),
                                "bbox": span.get("bbox", (0, 0, 0, 0)),
                                "color": span.get("color", 0),
                            })

                        if line_text.strip():
                            line_bbox = line.get("bbox", block_bbox)
                            page_spans.append({
                                "text": line_text.strip(),
                                "spans": line_spans,
                                "bbox": line_bbox,
                                "page_idx": page_idx,
                            })

                all_page_spans.append(page_spans)

            header_footer_texts = self._detect_headers_footers(all_page_spans, page_count, page_heights)

            font_sizes = []
            for page_spans in all_page_spans:
                for span_info in page_spans:
                    for s in span_info["spans"]:
                        font_sizes.append(s["size"])

            heading_thresholds = self._compute_heading_thresholds(font_sizes)

            md_lines = []
            total_text_len = 0

            bookmarks_text = self._extract_pdf_bookmarks(doc)
            if bookmarks_text:
                md_lines.append("\n## 目录\n\n")
                md_lines.append(bookmarks_text)
                md_lines.append("\n")

            metadata = self._extract_pdf_metadata(doc)
            if metadata:
                md_lines.append(self._format_pdf_metadata(metadata))
                md_lines.append("\n")

            for page_idx, page_spans in enumerate(all_page_spans):
                page = doc[page_idx]
                has_figure = False

                for span_info in page_spans:
                    text = span_info["text"]
                    bbox = span_info["bbox"]

                    if self._is_header_footer(text, bbox, header_footer_texts, page_dict_for_size=None):
                        continue

                    text = self._clean_text(text)
                    if not text:
                        continue

                    heading_level = self._detect_heading_level(span_info["spans"], heading_thresholds)

                    if heading_level:
                        prefix = "#" * heading_level
                        md_lines.append(f"\n{prefix} {text}\n")
                    else:
                        md_lines.append(f"{text}\n")

                    total_text_len += len(text)

                    if re.search(r'(Figure\s+\d+[:.])|(?<![如见看])(图\s*\d+)', text):
                        has_figure = True

                extracted_images_md = ""
                try:
                    image_list = page.get_images(full=True)
                    if image_list:
                        if not assets_dir_created:
                            assets_dir.mkdir(parents=True, exist_ok=True)
                            assets_dir_created = True

                        for img_idx, img_info in enumerate(image_list):
                            xref = img_info[0]
                            try:
                                base_image = doc.extract_image(xref)
                                if not base_image:
                                    continue

                                image_data = base_image.get("image", b"")
                                image_ext = base_image.get("ext", "png")
                                image_width = base_image.get("width", 0)
                                image_height = base_image.get("height", 0)
                                image_name = f"page_{page_idx+1}_img_{img_idx+1}.{image_ext}"

                                image_info = self._analyze_pdf_image(image_data, image_width, image_height)
                                if image_info.get("skip_reason"):
                                    self.logger.info(f"跳过PDF图片: {image_name} - {image_info['skip_reason']}")
                                    continue

                                image_path = assets_dir / image_name
                                with open(image_path, "wb") as fp:
                                    fp.write(image_data)

                                relative_path = f"{assets_dir_name}/{image_name}"
                                caption = self._generate_image_caption(image_info, page_idx + 1, img_idx + 1)

                                if image_info.get("is_chart"):
                                    extracted_images_md += f"\n> 📊 {caption}\n![{caption}]({relative_path})\n\n"
                                else:
                                    extracted_images_md += f"\n![{caption}]({relative_path})\n\n"

                                self.logger.info(f"提取图片: {image_name} ({image_info.get('size_str', 'unknown')})")
                            except Exception as img_e:
                                self.logger.warning(f"提取图片 xref={xref} 失败: {img_e}")
                            except Exception as img_e:
                                self.logger.warning(f"提取图片 xref={xref} 失败: {img_e}")
                except Exception as img_e:
                    self.logger.warning(f"提取第 {page_idx+1} 页图片时出错: {img_e}")

                if has_figure:
                    try:
                        self.logger.info(f"第 {page_idx+1} 页包含Figure，尝试渲染页面快照...")
                        if not assets_dir_created:
                            assets_dir.mkdir(parents=True, exist_ok=True)
                            assets_dir_created = True

                        pix = page.get_pixmap(dpi=150)
                        render_filename = f"page_{page_idx+1}_render.png"
                        render_path = assets_dir / render_filename
                        pix.save(str(render_path))

                        relative_path = f"{assets_dir_name}/{render_filename}"
                        extracted_images_md += f"\n> **Page {page_idx+1} Snapshot (Contains Figure)**\n\n![Page {page_idx+1} Render]({relative_path})\n\n"
                    except Exception as render_e:
                        self.logger.warning(f"渲染第 {page_idx+1} 页失败: {render_e}")

                if extracted_images_md:
                    md_lines.append(f"\n## 提取的图片\n\n{extracted_images_md}")

            md_text = "\n".join(md_lines)

            if total_text_len < 100 and pypdf_available:
                self.logger.info(f"{pdf_path.name} 可能是扫描版PDF，尝试使用OCR...")
                ocr_result = self._ocr_pdf(pdf_path)
                if ocr_result and len(ocr_result.strip()) > total_text_len:
                    md_text = ocr_result

            return md_text

        finally:
            doc.close()

    def _extract_pdf_bookmarks(self, doc) -> str:
        if not fitz_available:
            return ""

        try:
            toc = doc.get_toc(simple=False)
            if not toc:
                return ""

            self.logger.info(f"提取到 {len(toc)} 个书签/大纲")

            lines = []
            for item in toc:
                level = item[0]
                title = item[1]
                page_num = item[2] if len(item) > 2 else 0

                indent = "  " * (level - 1)
                page_ref = f" (第 {page_num} 页)" if page_num > 0 else ""
                title_clean = self._clean_markdown_title(title)
                lines.append(f"{indent}- [{title_clean}{page_ref}](#{self._generate_anchor(title_clean)})")

            return "\n".join(lines)

        except Exception as e:
            self.logger.warning(f"提取PDF书签失败: {e}")
            return ""

    def _clean_markdown_title(self, title: str) -> str:
        title = title.strip()
        title = re.sub(r'\s+', ' ', title)
        title = re.sub(r'[\[\]]', '', title)
        return title

    def _generate_anchor(self, title: str) -> str:
        anchor = title.lower()
        anchor = unicodedata.normalize('NFKC', anchor)
        parts = re.split(r'[\s\-_\.,;:!?()（）【】《》\[\]]+', anchor)
        anchor_parts = []
        for part in parts:
            part = part.strip()
            if not part:
                continue
            if re.match(r'^[\w]+$', part):
                anchor_parts.append(part)
            else:
                chinese_part = re.sub(r'[^\u4e00-\u9fff]', '', part)
                if chinese_part:
                    anchor_parts.append(chinese_part)
        anchor = '-'.join(anchor_parts)
        anchor = re.sub(r'-+', '-', anchor)
        anchor = anchor.strip('-')
        return anchor if anchor else title.lower().replace(' ', '-')

    def _extract_pdf_metadata(self, doc) -> Optional[Dict]:
        if not fitz_available:
            return None

        try:
            metadata = doc.metadata
            if not metadata:
                return None

            useful_fields = {}
            if metadata.get('title'):
                useful_fields['title'] = metadata['title']
            if metadata.get('author'):
                useful_fields['author'] = metadata['author']
            if metadata.get('subject'):
                useful_fields['subject'] = metadata['subject']
            if metadata.get('creator'):
                useful_fields['creator'] = metadata['creator']
            if metadata.get('producer'):
                useful_fields['producer'] = metadata['producer']
            if metadata.get('creationDate'):
                useful_fields['creation_date'] = metadata['creationDate']
            if metadata.get('modDate'):
                useful_fields['modification_date'] = metadata['modDate']

            return useful_fields if useful_fields else None

        except Exception as e:
            self.logger.warning(f"提取PDF元数据失败: {e}")
            return None

    def _format_pdf_metadata(self, metadata: Dict) -> str:
        lines = ["\n## 文档信息\n\n"]
        lines.append("| 属性 | 值 |\n")
        lines.append("|------|----|\n")

        field_names = {
            'title': '标题',
            'author': '作者',
            'subject': '主题',
            'creator': '创建程序',
            'producer': 'PDF生成器',
            'creation_date': '创建日期',
            'modification_date': '修改日期'
        }

        for key, value in metadata.items():
            display_name = field_names.get(key, key)
            lines.append(f"| {display_name} | {value} |\n")

        return "".join(lines)

    def _extract_pdf_with_pypdf(self, pdf_path: Path) -> str:
        if not pypdf_available:
            self.logger.error("pypdf也未安装，无法处理PDF文件")
            return ""

        try:
            reader = pypdf.PdfReader(pdf_path)
            text = ""

            metadata = {}
            if reader.metadata:
                metadata = {
                    k: v for k, v in reader.metadata.items()
                    if v and k in ['/Title', '/Author', '/Subject', '/Creator', '/Producer']
                }

            for page in reader.pages:
                page_text = page.extract_text()
                if page_text and page_text.strip():
                    page_text = self._clean_text(page_text)
                    text += page_text + "\n\n"

            if metadata:
                formatted_meta = self._format_pypdf_metadata(metadata)
                text = formatted_meta + "\n" + text

            if len(text.strip()) < 100:
                self.logger.info(f"{pdf_path.name} 可能是扫描版PDF，尝试使用OCR...")
                ocr_result = self._ocr_pdf(pdf_path)
                if ocr_result and len(ocr_result.strip()) > len(text.strip()):
                    text = ocr_result

            return self._convert_to_markdown(text) if text else ""

        except Exception as e:
            self.logger.error(f"处理 {pdf_path} 时出错: {str(e)}")
            return ""

    def _format_pypdf_metadata(self, metadata: Dict) -> str:
        lines = ["\n## 文档信息\n\n"]
        lines.append("| 属性 | 值 |\n")
        lines.append("|------|----|\n")

        field_names = {
            '/Title': '标题',
            '/Author': '作者',
            '/Subject': '主题',
            '/Creator': '创建程序',
            '/Producer': 'PDF生成器'
        }

        for key, value in metadata.items():
            display_name = field_names.get(key, key)
            lines.append(f"| {display_name} | {value} |\n")

        return "".join(lines)

    def _compute_heading_thresholds(self, font_sizes: List[float]) -> Dict[int, float]:
        if not font_sizes:
            return {}

        size_counts = Counter(font_sizes)
        body_size = size_counts.most_common(1)[0][0]

        unique_sizes = sorted(set(font_sizes), reverse=True)

        thresholds = {}
        heading_sizes = [s for s in unique_sizes if s > body_size]

        for i, size in enumerate(heading_sizes):
            level = i + 1
            if level <= 5:
                thresholds[level] = size

        return thresholds

    def _detect_heading_level(self, spans: List[Dict], heading_thresholds: Dict[int, float]) -> Optional[int]:
        if not spans or not heading_thresholds:
            return None

        max_size = max(s["size"] for s in spans)
        any_bold = any(self._is_bold_span(s) for s in spans)

        text = "".join(s["text"] for s in spans).strip()
        if len(text) > 200:
            return None

        if len(text.split()) < 2 and not any_bold:
            return None

        for level in sorted(heading_thresholds.keys()):
            if max_size >= heading_thresholds[level]:
                if level == 1 and not any_bold and max_size < heading_thresholds.get(1, 999) * 1.1:
                    pass
                return level

        if any_bold and max_size >= list(heading_thresholds.values())[-1] if heading_thresholds else False:
            return max(heading_thresholds.keys())

        return None

    def _is_bold_span(self, span: Dict) -> bool:
        font_name = span.get("font", "").lower()
        flags = span.get("flags", 0)
        if flags & 2 ** 4:
            return True
        bold_keywords = ["bold", "black", "heavy", "demi", "粗"]
        return any(kw in font_name for kw in bold_keywords)

    def _detect_headers_footers(self, all_page_spans: List[List[Dict]], page_count: int, page_heights: List[float] = None) -> Dict[str, int]:
        if page_count < 3:
            return {}

        text_position_map = Counter()

        for page_idx, page_spans in enumerate(all_page_spans):
            for span_info in page_spans:
                text = span_info["text"].strip()
                if len(text) > 100:
                    continue
                if re.match(r'^\d+$', text):
                    continue

                bbox = span_info["bbox"]
                y_pos = round(bbox[1], 0)
                key = f"{text}||{y_pos}"
                text_position_map[key] += 1

        header_footer_texts = {}
        threshold = max(3, page_count * 0.4)

        for key, count in text_position_map.items():
            if count >= threshold:
                text = key.split("||")[0]
                header_footer_texts[text] = count

        if page_heights:
            header_zone = set()
            footer_zone = set()
            for page_idx, page_spans in enumerate(all_page_spans):
                if page_idx >= len(page_heights):
                    break
                page_height = page_heights[page_idx]
                top_threshold = page_height * 0.1
                bottom_threshold = page_height * 0.9
                for span_info in page_spans:
                    bbox = span_info["bbox"]
                    y_top = bbox[1]
                    y_bottom = bbox[3]
                    text = span_info["text"].strip()
                    if not text or len(text) > 100:
                        continue
                    if y_top < top_threshold or y_bottom > bottom_threshold:
                        if text not in header_footer_texts:
                            header_footer_texts[text] = 1
                        else:
                            header_footer_texts[text] += 1

        return header_footer_texts

    def _is_header_footer(self, text: str, bbox: Tuple, header_footer_texts: Dict[str, int], page_dict_for_size=None) -> bool:
        if not header_footer_texts:
            return False

        stripped = text.strip()
        if stripped in header_footer_texts:
            return True

        for hf_text in header_footer_texts:
            if stripped == hf_text:
                return True

        return False

    def _analyze_pdf_image(self, image_data: bytes, width: int, height: int) -> Dict:
        info = {
            "width": width,
            "height": height,
            "size": len(image_data),
            "size_str": "",
            "is_chart": False,
            "is_icon": False,
            "skip_reason": None,
        }

        if info["size"] < 5 * 1024:
            info["skip_reason"] = f"图片过小 ({info['size']} bytes)"
            return info

        size_kb = info["size"] / 1024
        size_mb = info["size"] / (1024 * 1024)
        if size_mb >= 1:
            info["size_str"] = f"{size_mb:.1f} MB"
        else:
            info["size_str"] = f"{size_kb:.0f} KB"

        if width > 0 and height > 0:
            if width < 50 or height < 50:
                info["skip_reason"] = f"尺寸过小 ({width}x{height})"
                return info
            if height < 200 and (width / height) > 3.0:
                info["skip_reason"] = f"宽高比过大 ({width}x{height}), 可能是横幅"
                return info
            if height < 120 and (width / height) > 1.5:
                info["skip_reason"] = f"宽高比过大 ({width}x{height}), 可能是分隔线"
                return info

            ratio = width / height if height > 0 else 0
            area = width * height

            if area > 50000 and 0.3 < ratio < 3.0:
                info["is_chart"] = True

            if area < 10000 and (ratio > 3.0 or ratio < 0.3):
                info["is_icon"] = True

        return info

    def _generate_image_caption(self, image_info: Dict, page_num: int, img_num: int) -> str:
        if image_info.get("is_chart"):
            return f"图表 {page_num}-{img_num}"
        if image_info.get("is_icon"):
            return f"图标 {page_num}-{img_num}"
        return f"图片 {page_num}-{img_num}"

    def _clean_text(self, text: str) -> str:
        try:
            text = text.encode('utf-8', errors='ignore').decode('utf-8')
            text = ''.join(char for char in text if ord(char) >= 32 or char in '\n\t\r')
        except Exception:
            text = text.encode('ascii', errors='ignore').decode('ascii')
        return text

    def _detect_language_for_ocr(self, image) -> str:
        try:
            lang_config = pytesseract.get_languages(config='')
            has_chinese = any('chi_sim' in lang or 'chi_tra' in lang for lang in lang_config)

            if has_chinese:
                return 'chi_sim+eng'
            else:
                return 'eng'
        except Exception:
            return 'chi_sim+eng'

    def _ocr_pdf(self, pdf_path: Path) -> str:
        if not pypdf_available:
            self.logger.error("OCR回退库未安装(pypdf/pytesseract/pdf2image)")
            return ""

        try:
            if self.tesseract_cmd:
                pytesseract.pytesseract.tesseract_cmd = self.tesseract_cmd

            images = convert_from_path(pdf_path, poppler_path=self.poppler_path)
            text_parts = []
            detected_lang = None

            for i, image in enumerate(images):
                self.logger.info(f"正在OCR第 {i+1}/{len(images)} 页...")

                if detected_lang is None:
                    detected_lang = self._detect_language_for_ocr(image)
                    self.logger.info(f"检测到OCR语言: {detected_lang}")

                text = pytesseract.image_to_string(image, lang=detected_lang)
                text_parts.append(text)

            combined_text = "\n\n".join(text_parts)

            combined_text = self._merge_split_table_cells(combined_text)

            return combined_text

        except Exception as e:
            self.logger.error(f"OCR处理失败: {str(e)}")
            return ""

    def _merge_split_table_cells(self, text: str) -> str:
        lines = text.split('\n')
        if not lines:
            return text

        table_regions = self._detect_table_regions(lines)
        if not table_regions:
            return text

        merged_lines = []
        last_end = 0

        for start, end in table_regions:
            merged_lines.extend(lines[last_end:start])
            region_lines = lines[start:end]
            merged_region = self._merge_table_region(region_lines)
            merged_lines.extend(merged_region)
            last_end = end

        merged_lines.extend(lines[last_end:])
        return '\n'.join(merged_lines)

    def _detect_table_regions(self, lines: List[str]) -> List[Tuple[int, int]]:
        if len(lines) < 2:
            return []

        regions = []
        i = 0
        while i < len(lines) - 1:
            line = lines[i].strip()
            next_line = lines[i + 1].strip() if i + 1 < len(lines) else ""

            if self._is_definitive_table_row(line) and self._is_definitive_table_row(next_line):
                pipe_count = line.count('|')
                next_pipe_count = next_line.count('|')
                if abs(pipe_count - next_pipe_count) <= 1:
                    start = i
                    end = i + 2
                    while end < len(lines) and self._is_definitive_table_row(lines[end].strip()):
                        end += 1
                    if end - start >= 2:
                        regions.append((start, end))
                        i = end
                        continue

            i += 1

        return regions

    def _is_definitive_table_row(self, line: str) -> bool:
        line = line.strip()
        if not line:
            return False

        if line.startswith('|') and line.endswith('|'):
            pipe_count = line.count('|')
            if pipe_count < 3:
                return False

            cells = [c.strip() for c in line.split('|')[1:-1]]
            if not cells:
                return False

            cell_lengths = [len(c) for c in cells]
            avg_len = sum(cell_lengths) / len(cell_lengths) if cell_lengths else 0

            if avg_len > 60:
                return False

            non_empty = sum(1 for c in cells if c and not re.match(r'^[\s\-:]+$', c))
            if non_empty == 0:
                return False

            return True

        return False

    def _merge_table_region(self, region_lines: List[str]) -> List[str]:
        merged_lines = []
        i = 0

        while i < len(region_lines):
            line = region_lines[i].strip()

            if self._is_table_separator(line):
                merged_lines.append(region_lines[i])
                i += 1
                continue

            if not self._is_definitive_table_row(line):
                merged_lines.append(region_lines[i])
                i += 1
                continue

            merged_row = line
            j = i + 1

            while j < len(region_lines):
                next_line = region_lines[j].strip()
                if not self._is_definitive_table_row(next_line):
                    break

                if self._should_merge_table_rows(merged_row, next_line):
                    merged_row = self._merge_table_rows(merged_row, next_line)
                    j += 1
                else:
                    break

            merged_lines.append(merged_row)
            i = j

        return merged_lines

    def _is_table_separator(self, line: str) -> bool:
        separator_pattern = r'^\|[\s\-:]+\|[\s\-:]+\|'
        return bool(re.match(separator_pattern, line)) or line.strip() == '---' or line.strip() == '---'

    def _looks_like_table_row(self, line: str) -> bool:
        if not line.startswith('|') and not line.endswith('|'):
            return False
        pipe_count = line.count('|')
        return pipe_count >= 3

    def _should_merge_table_rows(self, row1: str, row2: str) -> bool:
        if not row1.startswith('|') or not row2.startswith('|'):
            return False

        cells1 = [c.strip() for c in row1.split('|')[1:-1]]
        cells2 = [c.strip() for c in row2.split('|')[1:-1]]

        if len(cells1) != len(cells2):
            return False

        non_empty_count1 = sum(1 for c in cells1 if c and not re.match(r'^[\s\-:]+$', c))
        non_empty_count2 = sum(1 for c in cells2 if c and not re.match(r'^[\s\-:]+$', c))

        return non_empty_count2 > non_empty_count1 * 0.5

    def _merge_table_rows(self, row1: str, row2: str) -> str:
        cells1 = [c.strip() for c in row1.split('|')[1:-1]]
        cells2 = [c.strip() for c in row2.split('|')[1:-1]]

        merged_cells = []
        for c1, c2 in zip(cells1, cells2):
            if c2 and not re.match(r'^[\s\-:]+$', c2):
                if c1:
                    merged_cells.append(f"{c1} {c2}")
                else:
                    merged_cells.append(c2)
            else:
                merged_cells.append(c1)

        return '|' + '|'.join(merged_cells) + '|'

    # ─────────────────────────────────────────────
    # 其他文档类型提取
    # ─────────────────────────────────────────────

    def _extract_text_from_word(self, docx_path: Path) -> str:
        if not docx_available:
            self.logger.error("docx2txt库未安装，无法处理Word文件")
            return "## Word内容提取失败\n\n需安装docx2txt库: pip install docx2txt"

        try:
            text = docx2txt.process(docx_path)
            return text

        except Exception as e:
            self.logger.error(f"处理Word文档 {docx_path} 时出错: {str(e)}")
            return ""

    def _extract_text_from_excel(self, excel_path: Path) -> str:
        if not pandas_available:
            self.logger.error("pandas/tabulate库未安装，无法处理Excel文件")
            return "## Excel内容提取失败\n\n需安装相关库: pip install pandas tabulate openpyxl"

        try:
            xls = pd.ExcelFile(excel_path)
            md_text = ""

            for sheet_name in xls.sheet_names:
                df = pd.read_excel(excel_path, sheet_name=sheet_name)
                md_text += f"## Sheet: {sheet_name}\n\n"
                md_table = df.to_markdown(index=False)
                md_text += md_table + "\n\n"

            return md_text

        except ImportError as e:
            if "openpyxl" in str(e):
                self.logger.error(f"处理Excel文档 {excel_path} 时出错: 缺少openpyxl依赖")
                return "## Excel内容提取失败\n\n需安装openpyxl库: pip install openpyxl\n\n错误信息: " + str(e)
            elif "tabulate" in str(e):
                self.logger.error(f"处理Excel文档 {excel_path} 时出错: 缺少tabulate依赖")
                return "## Excel内容提取失败\n\n需安装tabulate库: pip install tabulate\n\n错误信息: " + str(e)
            else:
                self.logger.error(f"处理Excel文档 {excel_path} 时出错: {str(e)}")
                return ""
        except Exception as e:
            self.logger.error(f"处理Excel文档 {excel_path} 时出错: {str(e)}")
            return ""

    def _extract_text_from_powerpoint(self, pptx_path: Path) -> str:
        if not pptx_available:
            self.logger.error("python-pptx库未安装，无法处理PowerPoint文件")
            return "## PowerPoint内容提取失败\n\n需安装python-pptx库: pip install python-pptx"

        try:
            prs = Presentation(pptx_path)
            md_text = ""

            for i, slide in enumerate(prs.slides):
                md_text += f"## 幻灯片 {i+1}\n\n"

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        md_text += f"{shape.text}\n\n"

                md_text += "---\n\n"

            return md_text

        except Exception as e:
            self.logger.error(f"处理PowerPoint文档 {pptx_path} 时出错: {str(e)}")
            return ""

    def _extract_text_from_html(self, html_path: Path) -> str:
        if not html2text_available:
            self.logger.error("html2text库未安装，无法处理HTML文件")
            return "## HTML内容提取失败\n\n需安装html2text库: pip install html2text"

        try:
            with open(html_path, "r", encoding="utf-8") as f:
                html_content = f.read()
            h = html2text.HTML2Text()
            h.ignore_links = False
            md_text = h.handle(html_content)
            return md_text
        except Exception as e:
            self.logger.error(f"处理HTML文档 {html_path} 时出错: {str(e)}")
            return ""

    def _convert_to_markdown(self, text: str) -> str:
        md_text = text

        md_text = re.sub(r'^\s*(\d+)\.\s+', r'\1. ', md_text, flags=re.MULTILINE)

        md_text = re.sub(r'\*([^*]+)\*', r'**\1**', md_text)

        return md_text

    def _optimize_markdown(self, md_text: str) -> str:
        md_text = self._preserve_tables(md_text)
        md_text = self._preserve_code_blocks(md_text)
        md_text = self._preserve_inline_code(md_text)
        md_text = self._preserve_links_and_images(md_text)
        md_text = self._preserve_math_formulas(md_text)
        md_text = self._preserve_blockquotes(md_text)
        md_text = self._detect_local_file_links(md_text)

        return md_text

    def _preserve_tables(self, text: str) -> str:
        lines = text.split('\n')
        result_lines = []
        i = 0

        while i < len(lines):
            line = lines[i]
            stripped = line.strip()

            if self._is_markdown_table_row(stripped):
                table_rows = [line]
                i += 1

                while i < len(lines):
                    next_line = lines[i]
                    next_stripped = next_line.strip()

                    if self._is_markdown_table_separator(next_stripped):
                        table_rows.append(next_line)
                        i += 1
                        break
                    elif self._is_markdown_table_row(next_stripped):
                        table_rows.append(next_line)
                        i += 1
                    else:
                        break

                while i < len(lines):
                    next_line = lines[i]
                    next_stripped = next_line.strip()

                    if self._is_markdown_table_row(next_stripped):
                        table_rows.append(next_line)
                        i += 1
                    else:
                        break

                for table_line in table_rows:
                    result_lines.append(table_line)
            else:
                result_lines.append(line)
                i += 1

        return '\n'.join(result_lines)

    def _is_markdown_table_row(self, line: str) -> bool:
        if not line.strip().startswith('|'):
            return False
        if '---' in line or ':--' in line or '--:' in line:
            return False
        return '|' in line[1:]

    def _is_markdown_table_separator(self, line: str) -> bool:
        cleaned = line.strip().strip('|').replace(' ', '')
        if not cleaned:
            return False
        segments = cleaned.split('|')
        return all(re.match(r'^:?-+:?$', seg) for seg in segments)

    def _preserve_code_blocks(self, text: str) -> str:
        lines = text.split('\n')
        result_lines = []
        i = 0
        in_code_block = False

        while i < len(lines):
            line = lines[i]

            if line.strip().startswith('```'):
                result_lines.append(line)
                in_code_block = not in_code_block
                i += 1
                continue

            if not in_code_block:
                if self._is_indented_code_block(line):
                    result_lines.append(line)
                else:
                    result_lines.append(line)
            else:
                result_lines.append(line)

            i += 1

        return '\n'.join(result_lines)

    def _is_indented_code_block(self, line: str) -> bool:
        return line.startswith('    ') or line.startswith('\t')

    def _preserve_inline_code(self, text: str) -> str:
        return text

    def _preserve_links_and_images(self, text: str) -> str:
        return text

    def _preserve_math_formulas(self, text: str) -> str:
        return text

    def _preserve_blockquotes(self, text: str) -> str:
        return text

    def _detect_local_file_links(self, text: str) -> str:
        pattern = r'\[([^\]]+)\]\(([^\)]+\.(?:pdf|docx?|xlsx?|pptx?|html?|txt))\)'
        matches = re.finditer(pattern, text)

        local_links = []
        for match in matches:
            link_text = match.group(1)
            link_path = match.group(2)

            if not link_path.startswith(('http://', 'https://', 'ftp://', 'mailto:')):
                local_links.append((link_text, link_path))

        if local_links:
            self.logger.info(f"检测到 {len(local_links)} 个本地文件链接:")
            for link_text, link_path in local_links:
                self.logger.info(f"  - [{link_text}]({link_path})")

        return text

    def _save_markdown(self, md_text: str, file_path: Path) -> str:
        output_path = Path(self.output_dir)
        output_path.mkdir(parents=True, exist_ok=True)

        file_name = file_path.stem
        file_type = file_path.suffix.lower()[1:]

        output_file = output_path / f"{file_name}.md"

        try:
            md_text = md_text.encode('utf-8', errors='ignore').decode('utf-8')
            md_text = ''.join(char for char in md_text if ord(char) >= 32 or char in '\n\t\r')
        except Exception as e:
            self.logger.warning(f"文本清理时出错: {e}，将使用替换策略")
            md_text = md_text.encode('ascii', errors='ignore').decode('ascii')

        header = f"""---
title: {file_name} Document
source_file: {file_path.name}
file_type: {file_type}
converted_date: {datetime.now().strftime("%Y-%m-%d %H:%M:%S")}
---

"""
        md_text = header + md_text

        try:
            with open(output_file, 'w', encoding='utf-8', errors='replace') as f:
                f.write(md_text)
        except Exception as e:
            self.logger.error(f"保存文件时出错: {e}")
            try:
                with open(output_file, 'w', encoding='ascii', errors='replace') as f:
                    f.write(md_text)
                self.logger.warning(f"使用ASCII编码保存文件: {output_file}")
            except Exception as e2:
                self.logger.error(f"使用ASCII编码保存文件也失败: {e2}")
                raise e2

        self.logger.info(f"已保存到 {output_file}")
        return str(output_file)
