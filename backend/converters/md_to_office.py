from typing import List, Optional
import os
from .base_converter import BaseConverter
from .batik_converter import BatikConverter
from .plantuml_converter import PlantUMLConverter
import json
import subprocess
import platform
import re
import shutil
from datetime import datetime
from pathlib import Path
import logging

# python-pptx相关导入（用于title_and_svg模式）
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from PIL import Image
    pptx_available = True
except ImportError:
    pptx_available = False

# 平台检测和条件导入（迁移自tools/md_to_docx.py）
IS_WINDOWS = platform.system() == "Windows"
if IS_WINDOWS:
    try:
        from win32com.client import Dispatch
        from docx.enum.section import WD_SECTION_START
        from docxcompose.composer import Composer
        from docx import Document
        from docxtpl import DocxTemplate
        WIN32COM_AVAILABLE = True
    except ImportError:
        WIN32COM_AVAILABLE = False
else:
    WIN32COM_AVAILABLE = False
    WD_SECTION_START = None

class MdToOfficeConverter(BaseConverter):
    """
    This class encapsulates the logic from the original md_to_docx.py script,
    refactored to be reusable and integrate into the VS Code extension backend.
    """

    def __init__(self, output_dir: str, **kwargs):
        super().__init__(output_dir, **kwargs)
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        
        # Get config from kwargs, with defaults
        self.output_format = kwargs.get('output_format', 'docx')
        self.docx_template_path = kwargs.get('docx_template_path')
        self.pptx_template_path = kwargs.get('pptx_template_path')
        self.project_name = kwargs.get('project_name', '')
        self.author = kwargs.get('author', '')
        self.mobilephone = kwargs.get('mobilephone', '')
        self.email = kwargs.get('email', '')
        self.promote_headings = kwargs.get('promote_headings', False)
        
        # PPTX SVG 转换模式配置 - 使用默认值
        self.pptx_svg_mode = 'full'
        
        # 模板路径现在由前端直接提供，后端不再进行复杂的查找
        self.template_path = None
        if self.output_format in ['docx', 'pdf']:
            self.template_path = self.docx_template_path
        elif self.output_format == 'pptx':
            self.template_path = self.pptx_template_path
        
        # 初始化Batik SVG转换器 - 将临时文件放到输出目录的svg_temp子目录
        svg_temp_dir = self.output_dir / 'svg_temp'
        self.batik_converter = BatikConverter(
            output_dir=str(svg_temp_dir),
            dpi=kwargs.get('svg_dpi', 300),
            timeout=kwargs.get('svg_timeout', 60)
        )

    def convert(self, input_path: str) -> List[str]:
        """
        Main conversion entry point. Handles both single files and directories.
        """
        if not self._is_valid_input(input_path, ['.md']):
            raise ValueError(f"Invalid input file or directory: {input_path}")

        output_files = []
        if os.path.isfile(input_path):
            output_file = self._convert_single_file(input_path)
            if output_file:
                output_files.append(output_file)
        else:
            md_files = self._get_files_by_extension(input_path, ['.md'])
            for md_file in md_files:
                output_file = self._convert_single_file(md_file)
                if output_file:
                    output_files.append(output_file)
        
        return output_files

    def _convert_single_file(self, input_file: str) -> Optional[str]:
        """
        Routes a single file to the correct conversion method based on output format.
        """
        if not Path(input_file).exists():
            self.logger.error(f"Input file not found: {input_file}")
            return None

        if self.output_format == 'docx':
            return self._convert_to_docx(input_file)
        elif self.output_format == 'pdf':
            # 1. 定义最终的PDF输出路径
            final_pdf_path = str(self.output_dir / f"{Path(input_file).stem}.pdf")
            
            # 2. 创建临时的DOCX文件
            docx_path = self._convert_to_docx(input_file, to_pdf=True)
            if not docx_path:
                self.logger.error(f"Failed to create intermediate DOCX for PDF conversion from {input_file}")
                return None
            
            # 3. 将临时DOCX转换为最终的PDF
            pdf_path_result = self._convert_docx_to_pdf(docx_path, final_pdf_path)
            
            # 4. 清理临时的DOCX文件
            if pdf_path_result and os.path.exists(docx_path):
                try:
                    os.remove(docx_path)
                    self.logger.info(f"Removed intermediate file: {docx_path}")
                except OSError as e:
                    self.logger.warning(f"Failed to remove intermediate file {docx_path}: {e}")

            return pdf_path_result
        elif self.output_format == 'html':
            return self._convert_to_html(input_file)
        elif self.output_format == 'pptx':
            return self._convert_to_pptx(input_file)
        else:
            self.logger.error(f"Unsupported output format: {self.output_format}")
            return None

    def _convert_to_pptx(self, input_file: str) -> Optional[str]:
        """Converts a Markdown file to PPTX."""
        input_path = Path(input_file)
        output_file_path = self.output_dir / f"{input_path.stem}.pptx"

        # 根据pptx_svg_mode选择不同的处理流程
        if self.pptx_svg_mode == 'title_and_svg':
            return self._process_title_and_svg_mode(input_file, output_file_path)
        else:
            return self._process_full_mode(input_file, output_file_path)

    def _process_title_and_svg_mode(self, input_file: str, output_file_path: Path) -> Optional[str]:
        """处理title_and_svg模式的PPTX转换"""
        if not pptx_available:
            self.logger.error("python-pptx库未安装，无法使用title_and_svg模式")
            return self._process_full_mode(input_file, output_file_path)
        
        self.logger.info(f"使用title_and_svg模式转换: {input_file}")
        
        try:
            # 读取并预处理Markdown内容
            processed_content, temp_files = self._preprocess_markdown(input_file)
            if processed_content is None:
                return None
            
            # 提取标题
            input_path = Path(input_file)
            title = self._get_title_from_md(processed_content, input_path)
            
            # 解析内容
            sections = self._parse_title_and_svg_mode(processed_content, title)
            
            # 创建演示文稿
            prs = self._create_presentation_from_template()
            
            # 获取Markdown文件所在目录
            md_dir = input_path.parent
            
            # 处理每个section
            for section in sections:
                if section['type'] == 'title':
                    self._create_title_slide(prs, section['title'])
                elif section['type'] == 'svg':
                    self._create_svg_slide(prs, section, md_dir)
            
            # 保存演示文稿
            prs.save(str(output_file_path))
            
            self.logger.info(f"Successfully converted {input_file} to {output_file_path}")
            return str(output_file_path)
            
        except Exception as e:
            self.logger.error(f"Failed during title_and_svg mode conversion: {e}")
            return None
        finally:
            # 清理临时文件
            for temp_file in temp_files:
                try:
                    if os.path.exists(temp_file):
                        os.remove(temp_file)
                except Exception as e:
                    self.logger.warning(f"Failed to remove temp file {temp_file}: {e}")
    
    def _process_full_mode(self, input_file: str, output_file_path: Path) -> Optional[str]:
        """处理full模式的PPTX转换（使用python-pptx创建多张幻灯片）"""
        if not pptx_available:
            self.logger.error("python-pptx库未安装，无法使用full模式")
            return None
        
        self.logger.info(f"使用full模式转换: {input_file}")
        
        try:
            # 读取并预处理Markdown内容
            processed_content, temp_files = self._preprocess_markdown(input_file)
            if processed_content is None:
                return None
            
            # 提取标题
            input_path = Path(input_file)
            title = self._get_title_from_md(processed_content, input_path)
            
            # 解析内容
            sections = self._parse_full_mode(processed_content, title)
            
            # 创建演示文稿
            prs = self._create_presentation_from_template()
            
            # 获取Markdown文件所在目录
            md_dir = input_path.parent
            
            # 处理每个section
            for section in sections:
                if section['type'] == 'content':
                    if section['level'] == 1 and not section['content']:
                        # 标题页
                        self._create_title_slide(prs, section['title'])
                    else:
                        # 内容页
                        self._create_content_slide(prs, section, md_dir)
            
            # 保存演示文稿
            prs.save(str(output_file_path))
            
            self.logger.info(f"Successfully converted {input_file} to {output_file_path}")
            return str(output_file_path)
            
        except Exception as e:
            self.logger.error(f"Failed during full mode conversion: {e}")
            return None
        finally:
            # 清理临时文件
            for temp_file in temp_files:
                try:
                    if os.path.exists(temp_file):
                        os.remove(temp_file)
                except Exception as e:
                    self.logger.warning(f"Failed to remove temp file {temp_file}: {e}")
    
    def _parse_title_and_svg_mode(self, content: str, title: str) -> List[dict]:
        """仅标题和SVG模式的解析（移植自MdToPptConverter）"""
        sections = []
        
        # 提取所有标题和图片的位置信息
        heading_pattern = re.compile(r'^(#+)\s+(.+)$', re.MULTILINE)
        headings = list(heading_pattern.finditer(content))
        
        # 提取所有SVG图片引用（包括已转换的PNG）
        svg_pattern = re.compile(r'!\[([^\]]*)\]\(([^\)]+\.(svg|png))\)', re.IGNORECASE)
        svg_matches = list(svg_pattern.finditer(content))
        
        # 添加文档标题作为第一页（只有当标题不为空且不与第一个标题重复时）
        first_heading_title = headings[0].group(2).strip() if headings else None
        if title and title != first_heading_title:
            sections.append({
                'level': 1, 
                'title': title, 
                'content': [],
                'type': 'title',
                'position': 0
            })
        
        # 创建包含标题和图片位置信息的列表
        items = []
        
        # 添加标题信息
        for match in headings:
            level = len(match.group(1))
            heading_text = match.group(2).strip()
            items.append({
                'type': 'title',
                'level': level,
                'title': heading_text,
                'position': match.start()
            })
        
        # 添加图片信息
        for match in svg_matches:
            alt_text = match.group(1)
            img_path = match.group(2)
            items.append({
                'type': 'svg',
                'level': 2,
                'title': alt_text or f"图片: {os.path.basename(img_path)}",
                'content': [match.group(0)],  # 完整的图片markdown语法
                'position': match.start()
            })
        
        # 按位置排序
        items.sort(key=lambda x: x['position'])
        
        # 将图片放在其前面最近的标题后面
        current_title_sections = []
        if title and title != first_heading_title:
            current_title_sections = [sections[0]]  # 文档标题
        
        for item in items:
            if item['type'] == 'title':
                # 添加标题页
                title_section = {
                    'level': item['level'],
                    'title': item['title'],
                    'content': [],
                    'type': 'title'
                }
                sections.append(title_section)
                current_title_sections.append(title_section)
            elif item['type'] == 'svg':
                # 图片放在最近的标题后面
                svg_section = {
                    'level': item['level'],
                    'title': item['title'],
                    'content': item['content'],
                    'type': 'svg'
                }
                sections.append(svg_section)
        
        return sections
    
    def _parse_full_mode(self, content: str, title: str) -> List[dict]:
        """完整模式的解析：标题占一页，内容根据情况分页，SVG单独占一页"""
        sections = []
        
        # 提取所有标题行
        heading_pattern = re.compile(r'^(#+)\s+(.+)$', re.MULTILINE)
        headings = list(heading_pattern.finditer(content))
        
        # 添加文档标题作为第一页（只有当标题不为空且不与第一个标题重复时）
        first_heading_title = headings[0].group(2).strip() if headings else None
        if title and title != first_heading_title:
            sections.append({
                'level': 1, 
                'title': title, 
                'content': [],
                'type': 'content'
            })
        
        # 处理每个标题和其内容
        for i, match in enumerate(headings):
            level = len(match.group(1))
            heading_text = match.group(2).strip()
            
            # 添加标题页
            sections.append({
                'level': level,
                'title': heading_text,
                'content': [],
                'type': 'content'
            })
            
            # 提取该标题下的内容
            content_start = match.end()
            if i + 1 < len(headings):
                content_end = headings[i + 1].start()
            else:
                content_end = len(content)
            
            section_content = content[content_start:content_end].strip()
            
            if section_content:
                # 将内容按段落分割
                content_lines = section_content.split('\n')
                
                # 过滤空行
                content_lines = [line for line in content_lines if line.strip()]
                
                if content_lines:
                    # 添加内容页
                    sections.append({
                        'level': level + 1,
                        'title': heading_text,
                        'content': content_lines,
                        'type': 'content'
                    })
        
        # 处理没有标题的情况（直接添加内容到第一节）
        if not headings and title:
            content_lines = content.strip().split('\n')
            content_lines = [line for line in content_lines if line.strip()]
            if content_lines:
                sections[0]['content'] = content_lines
        
        # 处理文件没有内容的情况
        if not sections:
            sections.append({'level': 1, 'title': title or '无标题', 'content': [], 'type': 'content'})
        
        return sections
    
    def _create_presentation_from_template(self) -> 'Presentation':
        """根据模板创建演示文稿，如果未提供模板则创建空白演示文稿"""
        if self.template_path and Path(self.template_path).exists():
            self.logger.info(f"正在加载模板: {self.template_path}")
            try:
                prs = Presentation(self.template_path)
                
                # 完全移除所有示例幻灯片，只保留布局
                while len(prs.slides) > 0:
                    rId = prs.slides._sldIdLst[-1].rId
                    prs.part.drop_rel(rId)
                    del prs.slides._sldIdLst[-1]

                self.logger.info(f"成功加载模板并移除所有示例幻灯片: {self.template_path}")
                return prs
            except Exception as e:
                self.logger.error(f"加载模板失败: {self.template_path}, 错误: {e}")
                self.logger.info("将创建空白演示文稿作为备用方案。")
                return Presentation()
        else:
            self.logger.info("未提供模板或模板不存在，正在创建空白演示文稿。")
            return Presentation()
    
    def _create_title_slide(self, prs: 'Presentation', title_text: str):
        """创建标题幻灯片"""
        # 使用占位符最少的布局
        min_placeholders = min(len(layout.placeholders) for layout in prs.slide_layouts)
        layout_to_use = None
        for layout in prs.slide_layouts:
            if len(layout.placeholders) == min_placeholders:
                layout_to_use = layout
                break
        
        slide = prs.slides.add_slide(layout_to_use)
        
        # 删除所有占位符形状
        shapes_to_remove = []
        for shape in slide.shapes:
            if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                shapes_to_remove.append(shape)
        
        for shape in shapes_to_remove:
            try:
                slide.shapes._spTree.remove(shape._element)
            except Exception as e:
                self.logger.warning(f"删除占位符时出错: {e}")
        
        # 设置白色背景
        try:
            # 方法1：通过slide.background设置
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(255, 255, 255)
        except Exception as e:
            try:
                # 方法2：通过添加白色矩形作为背景
                from pptx.enum.shapes import MSO_SHAPE
                bg_shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, 0, 0, 
                    slide.slide_layout.slide_master.slide_width,
                    slide.slide_layout.slide_master.slide_height
                )
                bg_fill = bg_shape.fill
                bg_fill.solid()
                bg_fill.fore_color.rgb = RGBColor(255, 255, 255)
                # 将背景形状移到最底层
                slide.shapes._spTree.insert(2, slide.shapes._spTree.pop())
                self.logger.info("使用矩形背景方法设置白色背景")
            except Exception as e2:
                self.logger.warning(f"设置背景颜色失败: 方法1={e}, 方法2={e2}")
        
        # 直接创建文本框，不使用占位符
        title_left = Inches(1)
        title_top = Inches(2.5)
        title_width = prs.slide_width - Inches(2)
        title_height = Inches(2)
        
        title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
        title_frame = title_box.text_frame
        title_frame.clear()
        
        p = title_frame.paragraphs[0]
        p.text = title_text
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.name = "微软雅黑"
    
    def _create_content_slide(self, prs: 'Presentation', section: dict, md_dir: Path):
        """创建内容幻灯片（包含标题、文本内容和图片）"""
        # 判断是否为纯标题页（没有内容）
        is_title_only = not section.get('content') or not any(content.strip() for content in section['content'])
        
        if is_title_only:
            # 纯标题页：使用标题页布局
            layout_to_use = None
            # 查找标题页布局
            for layout in prs.slide_layouts:
                name_lower = layout.name.lower()
                if any(keyword in name_lower for keyword in ['title', '标题', 'section']):
                    layout_to_use = layout
                    break
            
            # 如果没有找到标题页布局，使用占位符最少的布局
            if layout_to_use is None:
                min_placeholders = min(len(layout.placeholders) for layout in prs.slide_layouts)
                for layout in prs.slide_layouts:
                    if len(layout.placeholders) == min_placeholders:
                        layout_to_use = layout
                        break
        else:
            # 有内容的页面：优先使用内容页布局（Blank布局）
            layout_to_use = None
            
            # 首先查找Blank布局或其他内容页布局
            for layout in prs.slide_layouts:
                name_lower = layout.name.lower()
                if any(keyword in name_lower for keyword in ['blank', '空白', 'content', '内容']):
                    if layout_to_use is None or len(layout.placeholders) < len(layout_to_use.placeholders):
                        layout_to_use = layout
            
            # 如果没有找到合适的内容页布局，使用占位符最少的布局
            if layout_to_use is None:
                min_placeholders = min(len(layout.placeholders) for layout in prs.slide_layouts)
                for layout in prs.slide_layouts:
                    if len(layout.placeholders) == min_placeholders:
                        layout_to_use = layout
                        break
        
        slide = prs.slides.add_slide(layout_to_use)
        
        # 删除所有占位符形状
        shapes_to_remove = []
        for shape in slide.shapes:
            if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                shapes_to_remove.append(shape)
        
        for shape in shapes_to_remove:
            try:
                slide.shapes._spTree.remove(shape._element)
            except Exception as e:
                self.logger.warning(f"删除占位符时出错: {e}")
        
        try:
            # 根据是否为纯标题页决定标题位置和对齐方式
            if is_title_only:
                # 纯标题页：标题居中显示
                title_left = Inches(1)
                title_top = prs.slide_height / 2 - Inches(0.6)  # 垂直居中
                title_width = prs.slide_width - Inches(2)
                title_height = Inches(1.2)
                title_alignment = PP_ALIGN.CENTER
            else:
                # 有内容的页面：标题在顶部
                title_left = Inches(0.8)
                title_top = Inches(0.6)
                title_width = prs.slide_width - Inches(1.6)  # 左右各0.8英寸边距
                title_height = Inches(1.0)
                title_alignment = PP_ALIGN.LEFT
            
            title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
            title_frame = title_box.text_frame
            title_frame.clear()
            
            p = title_frame.paragraphs[0]
            p.text = section['title']
            p.alignment = title_alignment
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.name = "微软雅黑"
            
            # 处理内容（只有非纯标题页才处理内容）
            if not is_title_only and section['content']:
                content_text = '\n'.join(section['content'])
                
                # 检查是否包含图片
                img_matches = list(re.finditer(r'!\[[^\]]*\]\(([^\)]+)\)', content_text))
                
                if img_matches:
                    # 包含图片的情况：分别处理文本和图片
                    text_content = content_text
                    
                    # 移除图片引用，只保留文本
                    for match in reversed(img_matches):  # 从后往前删除，避免位置偏移
                        text_content = text_content[:match.start()] + text_content[match.end():]
                    
                    text_content = text_content.strip()
                    
                    # 添加文本内容（如果有）
                    if text_content:
                        # 在标题和内容之间添加红线分隔
                        line_left = Inches(0.8)
                        line_top = Inches(1.7)
                        line_width = prs.slide_width - Inches(1.6)
                        line_height = Inches(0.01)  # 1pt高度的细线
                        
                        from pptx.enum.shapes import MSO_SHAPE
                        red_line = slide.shapes.add_shape(
                            MSO_SHAPE.RECTANGLE, line_left, line_top, line_width, line_height
                        )
                        red_line_fill = red_line.fill
                        red_line_fill.solid()
                        red_line_fill.fore_color.rgb = RGBColor(200, 200, 200)  # 浅灰色
                        red_line.line.color.rgb = RGBColor(200, 200, 200)  # 浅灰色边框
                        
                        # 直接创建文本框，不使用占位符 - 增加安全边距
                        content_left = Inches(0.8)
                        content_top = Inches(1.9)  # 调整位置，在红线下方
                        content_width = prs.slide_width - Inches(1.6)  # 左右各0.8英寸边距
                        content_height = prs.slide_height - Inches(2.7)  # 调整高度，留出更多空间
                        
                        content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
                        content_frame = content_box.text_frame
                        content_frame.clear()
                        content_frame.word_wrap = True  # 启用自动换行
                        content_frame.auto_size = None  # 禁用自动调整大小
                        
                        content_p = content_frame.paragraphs[0]
                        content_p.text = text_content
                        content_p.alignment = PP_ALIGN.LEFT
                        content_p.font.size = Pt(18)
                        content_p.font.name = "微软雅黑"
                    
                    # 处理图片（每个图片单独占一页）
                    for match in img_matches:
                        img_path = match.group(1)
                        if not os.path.isabs(img_path):
                            img_path = os.path.join(md_dir, img_path)
                        
                        # 为每个图片创建单独的幻灯片
                        self._create_image_slide(prs, img_path, md_dir)
                        
                else:
                    # 纯文本内容，使用分页处理
                    self._add_text_with_pagination(prs, slide, content_text, section['title'])
                    
        except Exception as e:
            self.logger.error(f"创建内容幻灯片时出错: {e}")
            import traceback
            self.logger.error(traceback.format_exc())
    
    def _add_text_with_pagination(self, prs: 'Presentation', slide, text: str, title: str = ""):
        """添加文本并处理自动换行和分页"""
        try:
            # 在标题和内容之间添加红线分隔
            line_left = Inches(0.8)
            line_top = Inches(1.7)
            line_width = prs.slide_width - Inches(1.6)
            line_height = Inches(0.01)  # 1pt高度的细线
            
            from pptx.enum.shapes import MSO_SHAPE
            red_line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, line_left, line_top, line_width, line_height
            )
            red_line_fill = red_line.fill
            red_line_fill.solid()
            red_line_fill.fore_color.rgb = RGBColor(200, 200, 200)  # 浅灰色
            red_line.line.color.rgb = RGBColor(200, 200, 200)  # 浅灰色边框
            
            # 文本框配置 - 增加更安全的边距
            content_left = Inches(0.8)
            content_top = Inches(1.9)  # 调整位置，在红线下方
            content_width = prs.slide_width - Inches(1.6)  # 左右各0.8英寸边距
            content_height = prs.slide_height - Inches(2.7)  # 调整高度，留出更多空间
            
            # 创建文本框并启用自动换行
            content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
            content_frame = content_box.text_frame
            content_frame.clear()
            content_frame.word_wrap = True  # 启用自动换行
            content_frame.auto_size = None  # 禁用自动调整大小
            
            # 字体配置
            font_size = Pt(18)
            font_name = "微软雅黑"
            line_height = 1.2
            
            # 直接添加文本，让PowerPoint自动处理换行
            content_p = content_frame.paragraphs[0]
            content_p.text = text
            content_p.alignment = PP_ALIGN.LEFT
            content_p.font.size = font_size
            content_p.font.name = font_name
            content_p.line_spacing = line_height

                
        except Exception as e:
            self.logger.error(f"添加文本时出错: {e}")
            import traceback
            self.logger.error(traceback.format_exc())
    
    def _create_new_content_slide(self, prs: 'Presentation', title: str = ""):
        """创建新的内容幻灯片用于分页"""
        # 优先使用内容页布局（Blank布局），如果没有则使用占位符最少的布局
        layout_to_use = None
        
        # 首先查找Blank布局或其他内容页布局
        for layout in prs.slide_layouts:
            name_lower = layout.name.lower()
            if any(keyword in name_lower for keyword in ['blank', '空白', 'content', '内容']):
                if layout_to_use is None or len(layout.placeholders) < len(layout_to_use.placeholders):
                    layout_to_use = layout
        
        # 如果没有找到合适的内容页布局，使用占位符最少的布局
        if layout_to_use is None:
            min_placeholders = min(len(layout.placeholders) for layout in prs.slide_layouts)
            for layout in prs.slide_layouts:
                if len(layout.placeholders) == min_placeholders:
                    layout_to_use = layout
                    break
        
        slide = prs.slides.add_slide(layout_to_use)
        
        # 删除所有占位符形状
        shapes_to_remove = []
        for shape in slide.shapes:
            if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                shapes_to_remove.append(shape)
        
        for shape in shapes_to_remove:
            try:
                slide.shapes._spTree.remove(shape._element)
            except Exception as e:
                self.logger.warning(f"删除占位符时出错: {e}")
        
        # 设置白色背景
        try:
            # 方法1：通过slide.background设置
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(255, 255, 255)
        except Exception as e:
            try:
                # 方法2：通过添加白色矩形作为背景
                from pptx.enum.shapes import MSO_SHAPE
                bg_shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, 0, 0, 
                    prs.slide_width, prs.slide_height
                )
                bg_fill = bg_shape.fill
                bg_fill.solid()
                bg_fill.fore_color.rgb = RGBColor(255, 255, 255)
                # 将背景形状移到最底层
                slide.shapes._spTree.insert(2, slide.shapes._spTree.pop())
                self.logger.info("使用矩形背景方法设置白色背景")
            except Exception as e2:
                self.logger.warning(f"设置背景颜色失败: 方法1={e}, 方法2={e2}")
        
        # 如果有标题，添加标题 - 增加安全边距
        if title:
            title_left = Inches(0.8)
            title_top = Inches(0.6)
            title_width = prs.slide_width - Inches(1.6)  # 左右各0.8英寸边距
            title_height = Inches(1.0)
            
            title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
            title_frame = title_box.text_frame
            title_frame.clear()
            
            title_p = title_frame.paragraphs[0]
            title_p.text = title + " (续)"
            title_p.alignment = PP_ALIGN.LEFT  # 内容页标题左对齐
            title_p.font.size = Pt(32)
            title_p.font.bold = True
            title_p.font.name = "微软雅黑"
            
            # 在标题和内容之间添加红线分隔
            line_left = Inches(0.8)
            line_top = Inches(1.7)
            line_width = prs.slide_width - Inches(1.6)
            line_height = Inches(0.02)  # 2pt高度的红线
            
            from pptx.enum.shapes import MSO_SHAPE
            red_line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, line_left, line_top, line_width, line_height
            )
            red_line_fill = red_line.fill
            red_line_fill.solid()
            red_line_fill.fore_color.rgb = RGBColor(200, 200, 200)  # 浅灰色
            red_line.line.color.rgb = RGBColor(200, 200, 200)  # 浅灰色边框
        
        return slide
    
    def _create_image_slide(self, prs: 'Presentation', img_path: str, md_dir: Path):
        """创建单独的图片幻灯片"""
        # 优先使用内容页布局（Blank布局），如果没有则使用占位符最少的布局
        layout_to_use = None
        
        # 首先查找Blank布局或其他内容页布局
        for layout in prs.slide_layouts:
            name_lower = layout.name.lower()
            if any(keyword in name_lower for keyword in ['blank', '空白', 'content', '内容']):
                if layout_to_use is None or len(layout.placeholders) < len(layout_to_use.placeholders):
                    layout_to_use = layout
        
        # 如果没有找到合适的内容页布局，使用占位符最少的布局
        if layout_to_use is None:
            min_placeholders = min(len(layout.placeholders) for layout in prs.slide_layouts)
            for layout in prs.slide_layouts:
                if len(layout.placeholders) == min_placeholders:
                    layout_to_use = layout
                    break
        
        slide = prs.slides.add_slide(layout_to_use)
        
        # 删除所有占位符形状
        shapes_to_remove = []
        for shape in slide.shapes:
            if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                shapes_to_remove.append(shape)
        
        for shape in shapes_to_remove:
            try:
                slide.shapes._spTree.remove(shape._element)
            except Exception as e:
                self.logger.warning(f"删除占位符时出错: {e}")
        
        # 设置白色背景
        try:
            # 方法1：通过slide.background设置
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(255, 255, 255)
        except Exception as e:
            try:
                # 方法2：通过添加白色矩形作为背景
                from pptx.enum.shapes import MSO_SHAPE
                bg_shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, 0, 0, 
                    prs.slide_width, prs.slide_height
                )
                bg_fill = bg_shape.fill
                bg_fill.solid()
                bg_fill.fore_color.rgb = RGBColor(255, 255, 255)
                # 将背景形状移到最底层
                slide.shapes._spTree.insert(2, slide.shapes._spTree.pop())
                self.logger.info("使用矩形背景方法设置白色背景")
            except Exception as e2:
                self.logger.warning(f"设置背景颜色失败: 方法1={e}, 方法2={e2}")
        
        try:
            self.logger.info(f"处理图片: {img_path}")
            
            # 检查是否为SVG文件，如果是则需要转换
            if img_path.lower().endswith('.svg'):
                if not os.path.exists(img_path):
                    self.logger.warning(f"SVG文件未找到: {img_path}")
                    return
                
                # 使用Batik转换器直接转换SVG文件
                try:
                    svg_temp_dir = self.output_dir / 'svg_temp'
                    svg_temp_dir.mkdir(exist_ok=True)
                    
                    # 读取SVG内容
                    with open(img_path, 'r', encoding='utf-8') as f:
                        svg_content = f.read()
                    
                    # 生成PNG文件名
                    png_filename = f"{Path(img_path).stem}.png"
                    png_path = svg_temp_dir / png_filename
                    
                    # 转换SVG到PNG
                    success, message = self.batik_converter.convert_to_file(img_path, str(png_path))
                    if success and png_path.exists():
                        img_path = str(png_path)
                        self.logger.info(f"SVG转换成功: {img_path}")
                    else:
                        self.logger.warning(f"SVG转换失败，跳过图片: {img_path}")
                        return
                except Exception as e:
                    self.logger.warning(f"SVG转换过程中出错: {e}，跳过图片")
                    return
            else:
                # 非SVG文件，直接检查是否存在
                if not os.path.exists(img_path):
                    self.logger.warning(f"图片文件未找到: {img_path}")
                    return
            
            # 居中显示图片
            with Image.open(img_path) as img:
                aspect_ratio = img.width / img.height
                
                # 计算最大可用空间
                max_width = prs.slide_width - Inches(2)
                max_height = prs.slide_height - Inches(2.5)
                
                # 根据宽高比计算实际尺寸
                if aspect_ratio > max_width / max_height:
                    img_width = max_width
                    img_height = max_width / aspect_ratio
                else:
                    img_height = max_height
                    img_width = max_height * aspect_ratio
                
                # 计算居中位置
                img_left = (prs.slide_width - img_width) / 2
                img_top = (prs.slide_height - img_height) / 2
                
                # 添加图片
                slide.shapes.add_picture(img_path, img_left, img_top, width=img_width, height=img_height)
                self.logger.info(f"成功添加图片到幻灯片: {img_path}")
                
        except Exception as e:
            self.logger.error(f"添加图片时出错: {e}")
            import traceback
            self.logger.error(traceback.format_exc())

    def _create_svg_slide(self, prs: 'Presentation', section: dict, md_dir: Path):
        """创建SVG图片幻灯片（保留用于title_and_svg模式）"""
        # 优先使用内容页布局（Blank布局），如果没有则使用占位符最少的布局
        layout_to_use = None
        
        # 首先查找Blank布局或其他内容页布局
        for layout in prs.slide_layouts:
            name_lower = layout.name.lower()
            if any(keyword in name_lower for keyword in ['blank', '空白', 'content', '内容']):
                if layout_to_use is None or len(layout.placeholders) < len(layout_to_use.placeholders):
                    layout_to_use = layout
        
        # 如果没有找到合适的内容页布局，使用占位符最少的布局
        if layout_to_use is None:
            min_placeholders = min(len(layout.placeholders) for layout in prs.slide_layouts)
            for layout in prs.slide_layouts:
                if len(layout.placeholders) == min_placeholders:
                    layout_to_use = layout
                    break
        
        slide = prs.slides.add_slide(layout_to_use)
        
        # 删除所有占位符形状
        shapes_to_remove = []
        for shape in slide.shapes:
            if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                shapes_to_remove.append(shape)
        
        for shape in shapes_to_remove:
            try:
                slide.shapes._spTree.remove(shape._element)
            except Exception as e:
                self.logger.warning(f"删除占位符时出错: {e}")
        
        # 设置白色背景
        try:
            # 方法1：通过slide.background设置
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(255, 255, 255)
        except Exception as e:
            try:
                # 方法2：通过添加白色矩形作为背景
                from pptx.enum.shapes import MSO_SHAPE
                bg_shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, 0, 0, 
                    prs.slide_width, prs.slide_height
                )
                bg_fill = bg_shape.fill
                bg_fill.solid()
                bg_fill.fore_color.rgb = RGBColor(255, 255, 255)
                # 将背景形状移到最底层
                slide.shapes._spTree.insert(2, slide.shapes._spTree.pop())
                self.logger.info("使用矩形背景方法设置白色背景")
            except Exception as e2:
                self.logger.warning(f"设置背景颜色失败: 方法1={e}, 方法2={e2}")
        
        try:
            # 提取图片路径
            content = section['content'][0] if section['content'] else ''
            img_match = re.search(r'!\[[^\]]*\]\((.*?)\)', content)
            if not img_match:
                self.logger.warning("未找到图片引用")
                return
            
            img_path = img_match.group(1)
            if not os.path.isabs(img_path):
                img_path = os.path.join(md_dir, img_path)
            
            self.logger.info(f"处理图片: {img_path}")
            
            # 检查是否为SVG文件，如果是则需要转换
            if img_path.lower().endswith('.svg'):
                if not os.path.exists(img_path):
                    self.logger.warning(f"SVG文件未找到: {img_path}")
                    return
                
                # 使用Batik转换器直接转换SVG文件
                try:
                    svg_temp_dir = self.output_dir / 'svg_temp'
                    svg_temp_dir.mkdir(exist_ok=True)
                    
                    # 读取SVG内容
                    with open(img_path, 'r', encoding='utf-8') as f:
                        svg_content = f.read()
                    
                    # 生成PNG文件名
                    png_filename = f"{Path(img_path).stem}.png"
                    png_path = svg_temp_dir / png_filename
                    
                    # 转换SVG到PNG
                    success, message = self.batik_converter.convert_to_file(img_path, str(png_path))
                    if success and png_path.exists():
                        img_path = str(png_path)
                        self.logger.info(f"SVG转换成功: {img_path}")
                    else:
                        self.logger.warning(f"SVG转换失败，跳过图片: {img_path}")
                        return
                except Exception as e:
                    self.logger.warning(f"SVG转换过程中出错: {e}，跳过图片")
                    return
            else:
                # 非SVG文件，直接检查是否存在
                if not os.path.exists(img_path):
                    self.logger.warning(f"图片文件未找到: {img_path}")
                    return
            
            # 居中显示图片
            with Image.open(img_path) as img:
                aspect_ratio = img.width / img.height
                
                # 计算最大可用空间
                max_width = prs.slide_width - Inches(2)
                max_height = prs.slide_height - Inches(2.5)
                
                # 根据宽高比计算实际尺寸
                if aspect_ratio > max_width / max_height:
                    img_width = max_width
                    img_height = max_width / aspect_ratio
                else:
                    img_height = max_height
                    img_width = max_height * aspect_ratio
                
                # 计算居中位置
                img_left = (prs.slide_width - img_width) / 2
                img_top = (prs.slide_height - img_height) / 2
                
                # 添加图片
                slide.shapes.add_picture(img_path, img_left, img_top, width=img_width, height=img_height)
                self.logger.info(f"成功添加图片到幻灯片: {img_path}")
                
        except Exception as e:
            self.logger.error(f"添加SVG图片时出错: {e}")
            import traceback
            self.logger.error(traceback.format_exc())

    def _check_tool_availability(self, tool_name: str) -> bool:
        """Checks if an external tool is available in the system's PATH."""
        return shutil.which(tool_name) is not None

    def _extract_original_title(self, content: str) -> str:
        """从原始内容中提取标题（在任何处理之前）"""
        try:
            # 首先尝试从YAML front matter提取
            pandoc_title_match = re.search(r'^---\s*\ntitle:\s*(.+?)\n', content, re.DOTALL)
            if pandoc_title_match:
                return pandoc_title_match.group(1).strip()
            
            # 然后尝试提取第一个一级标题
            first_heading_match = re.search(r'^#\s+(.+)', content, re.MULTILINE)
            if first_heading_match:
                return first_heading_match.group(1).strip()
        except Exception as e:
            self.logger.warning(f"Could not extract original title due to error: {e}")
        
        return ""

    def _get_title_from_md(self, content: str, fallback_path: Path) -> str:
        """Extracts title from Markdown content."""
        # 如果有模板且已保存原始标题，使用原始标题
        has_template = self.template_path and Path(self.template_path).exists()
        if has_template and hasattr(self, '_original_title') and self._original_title:
            return self._original_title
            
        # 否则从当前内容中提取标题
        try:
            pandoc_title_match = re.search(r'^---\s*\ntitle:\s*(.+?)\n', content, re.DOTALL)
            if pandoc_title_match:
                return pandoc_title_match.group(1).strip()
            
            first_heading_match = re.search(r'^#\s+(.+)', content, re.MULTILINE)
            if first_heading_match:
                return first_heading_match.group(1).strip()
        except Exception as e:
            self.logger.warning(f"Could not extract title due to error: {e}")
        
        return fallback_path.stem

    def _preprocess_markdown(self, md_file_path: str) -> (Optional[str], List[str]):
        """
        Pre-processes Markdown content for conversion.
        Includes SVG processing and Mermaid diagram conversion.
        """
        try:
            with open(md_file_path, 'r', encoding='utf-8') as f:
                content = f.read()
        except Exception as e:
            self.logger.error(f"Cannot read Markdown file {md_file_path}: {e}")
            return None, []

        temp_files = []
        md_dir = Path(md_file_path).parent
        
        # 在标题提级之前先保存原始标题（用于模板中的{{title}}）
        self._original_title = self._extract_original_title(content)

        # 标题序号处理
        content = re.sub(r'^(#+)\s*(\d+(\.*\d+)*\s+)', r'\1 ', content, flags=re.MULTILINE)
        content = re.sub(r'^(#+)\s*(\d+(\.*\d+)*\.\s+)', r'\1 ', content, flags=re.MULTILINE)
        content = re.sub(r'(!\[)(fig:.*?)(\])', r'\1\3', content)
        
        # 自定义标题提级处理：二级标题提为一级，一级标题保持一级
        if self.promote_headings:
            content = self._custom_promote_headings(content)

        # SVG处理 - 使用Batik转换器
        try:
            self.logger.info(f"开始SVG处理，输出格式: {self.output_format}")
            
            # 所有格式都使用Batik转换器
            md_filename = Path(md_file_path).stem
            # 简化处理，直接使用原始内容
            processed_content = content
            svg_temp_files = []
            
            self.logger.info(f"SVG处理完成，使用Batik转换器")
        except Exception as e:
            self.logger.error(f"SVG处理失败: {e}")
            processed_content = content
            svg_temp_files = []
            stats = converter.get_conversion_statistics()
            if stats['svg_converted'] > 0:
                self.logger.info(f"SVG处理完成，转换了 {stats['svg_converted']} 个SVG块")
                self.logger.info(f"成功: {stats['svg_converted']}, 失败: {stats['conversion_failed']}")
                
                # 获取生成的文件列表
                if 'converted_files' in conversion_info:
                    temp_files.extend(conversion_info['converted_files'])
                if 'files_created' in stats:
                    temp_files.extend(stats['files_created'])
        except Exception as e:
            self.logger.error(f"SVG处理失败: {e}")
            # SVG处理失败不影响其他功能，继续执行

        # PlantUML文件链接处理
        try:
            self.logger.info("开始处理PlantUML文件链接...")
            content, plantuml_temp_files = self._process_plantuml_file_links(content, md_dir)
            temp_files.extend(plantuml_temp_files)
            self.logger.info(f"PlantUML文件链接处理完成，生成了 {len(plantuml_temp_files)} 个PNG文件")
        except Exception as e:
            self.logger.error(f"PlantUML文件链接处理失败: {e}")

        # Mermaid图表处理
        if self._check_tool_availability("mmdc"):
            def replace_mermaid(match):
                code = match.group(1)
                img_path = md_dir / f"mermaid-generated-{os.urandom(4).hex()}.png"
                try:
                    subprocess.run(['mmdc', '-i', '-', '-o', str(img_path)], input=code.encode('utf-8'), check=True, capture_output=True)
                    temp_files.append(str(img_path))
                    return f"![Mermaid Diagram]({img_path.name})"
                except (subprocess.CalledProcessError, FileNotFoundError) as e:
                    self.logger.error(f"Mermaid conversion failed: {e.stderr if hasattr(e, 'stderr') else e}")
                    return f"```mermaid\n{code}\n```"
            content = re.sub(r'```mermaid\n(.*?)\n```', replace_mermaid, content, flags=re.DOTALL)

        # 表格列宽优化处理
        content = self._optimize_table_column_widths(content)

        return content, temp_files
    
    def _custom_promote_headings(self, content: str) -> str:
        """
        自定义标题提级处理：
        - 有模板时：一级标题作为{{title}}变量，不出现在正文中；二级标题提升为一级标题
        - 无模板时：一级标题转为不带序号的大字体正文；二级标题提升为一级标题
        - 三级及以下标题相应提升一级
        """
        lines = content.split('\n')
        processed_lines = []
        has_template = self.template_path and Path(self.template_path).exists()
        
        for line in lines:
            # 匹配标题行
            heading_match = re.match(r'^(#+)\s+(.+)$', line)
            if heading_match:
                heading_level = len(heading_match.group(1))
                heading_text = heading_match.group(2)
                
                if heading_level == 1:
                    if has_template:
                        # 有模板时：一级标题不出现在正文中（作为{{title}}变量使用）
                        continue
                    else:
                        # 无模板时：一级标题转为不带序号的大字体正文
                        processed_lines.append(f'**{heading_text}**\n')
                elif heading_level == 2:
                    # 二级标题提升为一级标题
                    processed_lines.append(f'# {heading_text}')
                elif heading_level >= 3:
                    # 三级及以下标题提升一级
                    new_level = heading_level - 1
                    processed_lines.append('#' * new_level + f' {heading_text}')
                else:
                    processed_lines.append(line)
            else:
                processed_lines.append(line)
        
        return '\n'.join(processed_lines)
    
    def _cleanup_temp_files(self, temp_files: List[str], processed_file: str = None, original_file: str = None, preserve_png_for_html: bool = False):
        """清理临时文件"""
        for temp_file in temp_files:
            try:
                if os.path.exists(temp_file):
                    # HTML转换时保留PNG文件，删除SVG文件
                    if preserve_png_for_html:
                        if temp_file.lower().endswith('.png'):
                            self.logger.info(f"HTML转换：保留PNG文件 {temp_file}")
                            continue
                        elif temp_file.lower().endswith('.svg'):
                            self.logger.info(f"HTML转换：删除SVG文件 {temp_file}")
                    
                    os.remove(temp_file)
            except Exception as e:
                self.logger.warning(f"无法删除临时文件 {temp_file}: {e}")
        
        # 清理处理过的文件（如果与原文件不同）
        if processed_file and original_file and processed_file != original_file:
            try:
                if os.path.exists(processed_file):
                    os.remove(processed_file)
            except Exception as e:
                self.logger.warning(f"无法删除临时文件 {processed_file}: {e}")
        
        # 清理svg_temp目录（如果存在且不是HTML转换）
        if not preserve_png_for_html:
            try:
                svg_temp_dir = self.output_dir / 'svg_temp'
                if svg_temp_dir.exists() and svg_temp_dir.is_dir():
                    shutil.rmtree(svg_temp_dir)
                    self.logger.info(f"已删除svg_temp目录: {svg_temp_dir}")
            except Exception as e:
                self.logger.warning(f"无法删除svg_temp目录: {e}")
        
        # 清理临时文件
        try:
            # 清理Batik转换器的临时文件（如果有的话）
            pass
        except Exception as e:
            self.logger.warning(f"Batik转换器临时文件清理失败: {e}")

    def _convert_to_docx(self, input_file: str, to_pdf: bool = False) -> Optional[str]:
        """Converts a Markdown file to DOCX using a unified pandoc approach."""
        input_path = Path(input_file)
        
        processed_content, temp_images = self._preprocess_markdown(input_file)
        if processed_content is None:
            return None

        processed_md_file = input_path.with_name(f"{input_path.stem}_processed_{os.getpid()}.md")
        processed_md_file.write_text(processed_content, encoding='utf-8')
        
        all_temp_files = temp_images + [str(processed_md_file)]

        try:
            if not self._check_tool_availability("pandoc"):
                self.logger.error("Pandoc not found. Please install pandoc and add it to your PATH.")
                raise FileNotFoundError("Pandoc not found. Please install pandoc and add it to your system's PATH.")

            # Decide whether to use the advanced template feature
            use_advanced_template = (
                self.template_path and 
                Path(self.template_path).exists() and 
                WIN32COM_AVAILABLE
            )

            if use_advanced_template:
                # --- Advanced Template Path (Windows Only) ---
                self.logger.info(f"使用高级模板功能: {self.template_path}")
                
                # 1. Create a temporary content-only DOCX
                temp_content_docx = self.output_dir / f"{input_path.stem}_content_{os.getpid()}.docx"
                all_temp_files.append(str(temp_content_docx))
                
                cmd = [
                    'pandoc', str(processed_md_file),
                    '-o', str(temp_content_docx),
                    '--resource-path=' + str(input_path.parent),
                    '--quiet'
                ]
                # 标题提级现在在预处理阶段处理，不再使用Pandoc参数
                # if self.promote_headings:
                #     cmd.append('--shift-heading-level-by=-1')
                    
                subprocess.run(cmd, check=True, capture_output=True, text=True, encoding='utf-8')

                # 2. Get title and compose final document
                title = self._get_title_from_md(processed_content, input_path)
                final_output_path = self._copy_template_and_append_content(
                    self.template_path,
                    str(temp_content_docx),
                    title,
                    original_input_file=input_file
                )

                # 3. Update TOC if composition was successful
                if final_output_path and Path(final_output_path).exists() and final_output_path != str(temp_content_docx):
                    self._update_toc(final_output_path)
                    self.logger.info(f"成功转换并应用模板: {input_file} -> {final_output_path}")
                    
                    # If converting to PDF, this is the intermediate file.
                    if to_pdf:
                        return final_output_path
                    
                    # Otherwise, it's the final product, we can clean up the content docx
                    # Note: all_temp_files will be cleaned up in the finally block.
                    return final_output_path
                else:
                    self.logger.warning("模板合成失败，返回无模板的DOCX文件。")
                    # It failed and returned the original content_path. We must rename it to a non-temp name.
                    final_path_on_failure = self.output_dir / f"{input_path.stem}.docx"
                    shutil.move(temp_content_docx, final_path_on_failure)
                    all_temp_files.remove(str(temp_content_docx)) # Don't delete it
                    return str(final_path_on_failure)
            else:
                # --- Simple/Cross-Platform Path ---
                if self.template_path:
                    if not WIN32COM_AVAILABLE:
                        self.logger.warning("检测到非Windows环境，模板功能受限（仅应用样式）。")
                    self.logger.info(f"使用DOCX模板进行样式转换 (reference-doc): {self.template_path}")
                else:
                    self.logger.info("未提供DOCX模板，使用Pandoc默认样式")

                if to_pdf:
                    output_file_path = self.output_dir / f"{input_path.stem}_temp_for_pdf_{os.getpid()}.docx"
                else:
                    output_file_path = self.output_dir / f"{input_path.stem}.docx"
                
                cmd = [
                    'pandoc', str(processed_md_file),
                    '-o', str(output_file_path),
                    '--resource-path=' + str(input_path.parent),
                    '--quiet'
                ]
                # Use --reference-doc for styling, same as PPTX conversion
                if self.template_path and Path(self.template_path).exists():
                    cmd.extend(['--reference-doc', self.template_path])

                # 标题提级现在在预处理阶段处理，不再使用Pandoc参数
                # if self.promote_headings:
                #     cmd.append('--shift-heading-level-by=-1')
                
                subprocess.run(cmd, check=True, capture_output=True, text=True, encoding='utf-8')
                
                self.logger.info(f"成功转换 {input_file} to {output_file_path}")
                return str(output_file_path)

        except (subprocess.CalledProcessError, FileNotFoundError) as e:
            error_message = e.stderr if hasattr(e, 'stderr') else str(e)
            self.logger.error(f"Failed during DOCX conversion: {error_message}")
            return None
        finally:
            self._cleanup_temp_files(all_temp_files)

    def _update_toc(self, docx_path: str):
        """Updates the Table of Contents in a DOCX file using Word COM object."""
        if not WIN32COM_AVAILABLE:
            return
        
        word = None
        try:
            word = Dispatch('Word.Application')
            doc = word.Documents.Open(str(Path(docx_path).resolve()))
            doc.Fields.Update()
            if hasattr(doc, 'TablesOfContents'):
                for toc in doc.TablesOfContents:
                    toc.Update()
            doc.Save()
            self.logger.info(f"Updated TOC for {docx_path}")
        except Exception as e:
            self.logger.error(f"Failed to update TOC for {docx_path}: {e}")
        finally:
            if word:
                try:
                    if 'doc' in locals() and doc:
                        doc.Close(False)
                    word.Quit()
                except:
                    pass
    
    def _convert_docx_to_pdf(self, docx_path: str, pdf_path: str) -> Optional[str]:
        """Converts a DOCX file to PDF."""
        final_pdf_path = Path(pdf_path)

        if WIN32COM_AVAILABLE:
            word = None
            try:
                word = Dispatch('Word.Application')
                doc = word.Documents.Open(str(Path(docx_path).resolve()))
                doc.SaveAs(str(final_pdf_path.resolve()), FileFormat=17)
                self.logger.info(f"Successfully created PDF with Word: {final_pdf_path}")
                return str(final_pdf_path)
            except Exception as e:
                self.logger.error(f"Word PDF conversion failed: {e}")
            finally:
                if word:
                    try:
                        if 'doc' in locals() and doc:
                            doc.Close(False)
                        word.Quit()
                    except:
                        pass

        if self._check_tool_availability("soffice"):
            try:
                # soffice 会自动处理输出文件名，我们只需提供目录
                cmd = ["soffice", "--headless", "--convert-to", "pdf", "--outdir", str(final_pdf_path.parent), docx_path]
                subprocess.run(cmd, check=True, capture_output=True)
                
                # LibreOffice/soffice 会创建与输入文件同名的PDF，但可能与我们期望的命名不同，所以需要重命名
                expected_soffice_output = Path(docx_path).with_suffix('.pdf')
                if expected_soffice_output.exists() and str(expected_soffice_output) != str(final_pdf_path):
                    shutil.move(str(expected_soffice_output), str(final_pdf_path))

                self.logger.info(f"Successfully created PDF with LibreOffice: {final_pdf_path}")
                return str(final_pdf_path)
            except (subprocess.CalledProcessError, FileNotFoundError) as e:
                self.logger.warning(f"LibreOffice conversion failed: {e.stderr if hasattr(e, 'stderr') else e}")

        self.logger.error("No suitable tool (Word/LibreOffice) found for PDF conversion.")
        return None

    def _convert_to_html(self, input_file: str) -> Optional[str]:
        """Converts a Markdown file to a styled HTML file."""
        input_path = Path(input_file)
        output_file_path = self.output_dir / f"{input_path.stem}.html"

        processed_content, temp_images = self._preprocess_markdown(input_file)
        if processed_content is None:
            return None

        processed_md_file = input_path.with_name(f"{input_path.stem}_processed_{os.getpid()}.md")
        processed_md_file.write_text(processed_content, encoding='utf-8')
        
        all_temp_files = temp_images + [str(processed_md_file)]

        try:
            if not self._check_tool_availability("pandoc"):
                self.logger.error("Pandoc not found. Please install it to convert files.")
                return None
            
            resource_path_arg = '--resource-path=' + str(input_path.parent)
            cmd = ['pandoc', str(processed_md_file), '--from', 'markdown+smart', '--to', 'html', resource_path_arg]
            result = subprocess.run(cmd, check=True, capture_output=True, text=True, encoding='utf-8')
            html_body = result.stdout
            
            heading_counts = {}
            def add_anchor_to_heading(match):
                level, title = len(match.group(1)), match.group(2).strip()
                base_id = re.sub(r'[^\w\s-]', '', title).strip().lower()
                base_id = re.sub(r'[\s-]+', '-', base_id)
                count = heading_counts.get(base_id, 0)
                heading_counts[base_id] = count + 1
                anchor_id = f"{base_id}-{count}" if count > 0 else base_id
                return f'<h{level} id="{anchor_id}">{title}</h{level}>'
            html_body = re.sub(r'<h([1-6])>(.*?)</h\1>', add_anchor_to_heading, html_body)

            toc_html = self._generate_html_toc(processed_content)
            title = self._get_title_from_md(processed_content, input_path)
            css = self._get_html_theme_css("github_floating_toc")
            
            final_html = f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{title}</title>
    <style>{css}</style>
</head>
<body>
    <div class="container">
        <div class="toc-container">{toc_html}</div>
        <div class="content-container">{html_body}</div>
    </div>
</body>
</html>"""

            output_file_path.write_text(final_html, encoding='utf-8')
            self.logger.info(f"Successfully created HTML: {output_file_path}")
            return str(output_file_path)

        except (subprocess.CalledProcessError, FileNotFoundError) as e:
            self.logger.error(f"Failed during HTML conversion: {e.stderr if hasattr(e, 'stderr') else e}")
            return None
        finally:
            # HTML转换时不清理PNG文件，因为它们需要保留在svg_temp目录中供HTML引用
            # 只清理处理过的markdown文件
            self._cleanup_temp_files([str(processed_md_file)], str(processed_md_file), input_file, preserve_png_for_html=True)

    def _generate_html_toc(self, content: str) -> str:
        """Generates a nested HTML list for the Table of Contents."""
        toc_lines = ['<nav class="toc"><ul>']
        heading_counts = {}
        for line in content.splitlines():
            match = re.match(r'^(#+)\s+(.*)', line)
            if match:
                level, title = len(match.group(1)), match.group(2).strip()
                base_id = re.sub(r'[^\w\s-]', '', title).strip().lower()
                base_id = re.sub(r'[\s-]+', '-', base_id)
                count = heading_counts.get(base_id, 0)
                heading_counts[base_id] = count + 1
                anchor_id = f"{base_id}-{count}" if count > 0 else base_id
                toc_lines.append(f'<li class="toc-level-{level}"><a href="#{anchor_id}">{title}</a></li>')
        toc_lines.append('</ul></nav>')
        return '\n'.join(toc_lines)

    def _get_html_theme_css(self, theme_name: str) -> str:
        """Returns CSS for the HTML output."""
        github_floating_toc = """
        body { font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Helvetica, Arial, sans-serif; line-height: 1.6; color: #333; background-color: #fff; margin: 0; padding: 0; }
        .container { max-width: 1200px; margin: 20px auto; display: flex; flex-direction: row; align-items: flex-start; }
        .toc-container { width: 250px; flex-shrink: 0; position: -webkit-sticky; position: sticky; top: 20px; height: calc(100vh - 40px); overflow-y: auto; padding-right: 20px; border-right: 1px solid #e1e4e8; }
        .content-container { flex-grow: 1; padding-left: 30px; max-width: 800px; }
        .toc ul { list-style: none; padding-left: 0; } .toc li a { color: #0366d6; text-decoration: none; display: block; padding: 4px 0; font-size: 14px; }
        .toc li a:hover { text-decoration: underline; }
        .toc-level-1 { padding-left: 5px; font-weight: 600; } .toc-level-2 { padding-left: 20px; } .toc-level-3 { padding-left: 35px; } .toc-level-4 { padding-left: 50px; }
        h1, h2, h3, h4, h5, h6 { font-weight: 600; line-height: 1.25; margin-top: 24px; margin-bottom: 16px; border-bottom: 1px solid #eaecef; padding-bottom: .3em; }
        h1 { font-size: 2em; } h2 { font-size: 1.5em; } h3 { font-size: 1.25em; }
        p { margin-top: 0; margin-bottom: 16px; } a { color: #0366d6; text-decoration: none; } a:hover { text-decoration: underline; }
        code, pre { font-family: "SFMono-Regular", Consolas, "Liberation Mono", Menlo, Courier, monospace; font-size: 13px; }
        pre { word-wrap: normal; padding: 16px; overflow: auto; line-height: 1.45; background-color: #f6f8fa; border-radius: 3px; }
        code { background-color: rgba(27,31,35,.05); padding: .2em .4em; margin: 0; border-radius: 3px; }
        pre > code { padding: 0; margin: 0; background-color: transparent; border: 0; }
        table { border-collapse: collapse; } th, td { border: 1px solid #ddd; padding: 8px; } th { background-color: #f2f2f2; }
        img { max-width: 100%; } blockquote { color: #6a737d; border-left: .25em solid #dfe2e5; padding: 0 1em; margin-left: 0; }
        """
        return github_floating_toc

    def _remove_title_numbers(self, input_file: str) -> str:
        """
        处理Markdown文件，去掉标题前面的序号（如1.1、2.3.4等格式），
        并删除图片标题，同时确保列表前有空行以便Pandoc正确识别。
        """
        try:
            # 读取原始文件内容
            with open(input_file, 'r', encoding='utf-8') as f:
                original_content = f.read()
            
            # 1. 正则表达式匹配标题前的序号
            pattern_title_numbers = r'^(#+)\s+(\d+(\.\d+)*)\s+(.+)$'
            processed_content = re.sub(pattern_title_numbers, r'\1 \4', original_content, flags=re.MULTILINE)
            
            # 2. 删除图片标题
            processed_content = self._remove_image_captions(processed_content)

            # 3. 确保列表前有空行以便Pandoc正确识别
            lines = processed_content.splitlines()
            new_processed_lines = []
            list_markers = ("- ", "* ", "+ ")
            
            for i, current_line_text in enumerate(lines):
                stripped_line = current_line_text.lstrip()
                is_list_item = any(stripped_line.startswith(marker) for marker in list_markers)

                if is_list_item:
                    if i > 0:
                        previous_line_text = lines[i-1]
                        stripped_previous_line = previous_line_text.lstrip()
                        is_previous_list_item = any(stripped_previous_line.startswith(marker) for marker in list_markers)
                        
                        if previous_line_text.strip() and not is_previous_list_item:
                            new_processed_lines.append("")
                
                new_processed_lines.append(current_line_text)
            
            processed_content = "\n".join(new_processed_lines)
            
            # 确保处理后的内容以换行符结尾
            if original_content.endswith('\n') and not processed_content.endswith('\n'):
                processed_content += '\n'
            elif new_processed_lines and any(new_processed_lines[-1].lstrip().startswith(marker) for marker in list_markers) and not processed_content.endswith('\n'):
                processed_content += '\n'

            # 如果内容没有变化，直接返回原文件路径
            if processed_content == original_content:
                return input_file
            
            # 创建临时文件保存处理后的内容
            temp_file = input_file + '.temp.md'
            with open(temp_file, 'w', encoding='utf-8') as f:
                f.write(processed_content)
            
            return temp_file
        except Exception as e:
            self.logger.error(f"处理Markdown预处理时出错 (file: {input_file}): {str(e)}")
            return input_file
    
    def _remove_image_captions(self, content: str) -> str:
        """删除图片标题（占位符实现）"""
        # 这里可以添加具体的图片标题删除逻辑
        return content
    
    def _copy_template_and_append_content(self, template_path: str, content_path: str, title: str, original_input_file: str) -> str:
        """
        Applies a template by rendering variables and composing it with the content document.
        This uses `DocxTemplate` and `docxcompose`.
        """
        if not WIN32COM_AVAILABLE:
            self.logger.warning("在非Windows系统上无法使用模板功能，将使用简单转换")
            return content
    
    def _process_plantuml_file_links(self, content: str, md_dir: Path) -> tuple[str, List[str]]:
        """
        处理Markdown中的PlantUML文件链接，将其转换为PNG图片链接
        
        Args:
            content: Markdown内容
            md_dir: Markdown文件所在目录
            
        Returns:
            tuple: (处理后的内容, 生成的临时文件列表)
        """
        temp_files = []
        
        # 匹配PlantUML文件链接的正则表达式
        # 支持 ![alt](path.puml), ![alt](path.plantuml), ![alt](path.pu)
        plantuml_pattern = re.compile(r'!\[([^\]]*)\]\(([^)]+\.(?:puml|plantuml|pu))\)', re.IGNORECASE)
        
        def replace_plantuml_link(match):
            alt_text = match.group(1)
            puml_path = match.group(2)
            
            # 处理相对路径
            if not os.path.isabs(puml_path):
                full_puml_path = md_dir / puml_path
            else:
                full_puml_path = Path(puml_path)
            
            # 检查PlantUML文件是否存在
            if not full_puml_path.exists():
                self.logger.warning(f"PlantUML文件不存在: {full_puml_path}")
                return match.group(0)  # 返回原始链接
            
            try:
                # 首先检查是否已经存在对应的PNG文件
                puml_stem = full_puml_path.stem
                existing_png_candidates = [
                    self.output_dir / f"{puml_stem}.png",
                    md_dir / f"{puml_stem}.png",
                    Path(f"{puml_stem}.png")
                ]
                
                existing_png = None
                for candidate in existing_png_candidates:
                    if candidate.exists():
                        existing_png = candidate
                        break
                
                if existing_png:
                    # 使用已存在的PNG文件
                    try:
                        relative_png_path = existing_png.relative_to(md_dir)
                    except ValueError:
                        # 如果无法计算相对路径，复制文件到Markdown目录
                        import shutil
                        target_png = md_dir / f"{puml_stem}.png"
                        shutil.copy2(existing_png, target_png)
                        relative_png_path = target_png.name
                        temp_files.append(str(target_png))
                    
                    new_link = f"![{alt_text}]({relative_png_path})"
                    self.logger.info(f"使用已存在的PNG文件: {puml_path} -> {relative_png_path}")
                    return new_link
                
                # 如果没有现成的PNG文件，尝试转换
                plantuml_converter = PlantUMLConverter(str(self.output_dir))
                
                # 转换PlantUML文件为PNG
                result = plantuml_converter.convert(str(full_puml_path))
                
                if result and len(result) > 0:
                    png_path = Path(result[0])
                    if png_path.exists():
                        # 记录临时文件
                        temp_files.append(str(png_path))
                        
                        # 计算相对于Markdown文件的PNG路径
                        try:
                            relative_png_path = png_path.relative_to(md_dir)
                        except ValueError:
                            # 如果无法计算相对路径，复制文件到Markdown目录
                            import shutil
                            target_png = md_dir / f"{puml_stem}.png"
                            shutil.copy2(png_path, target_png)
                            relative_png_path = target_png.name
                            temp_files.append(str(target_png))
                        
                        # 返回新的图片链接
                        new_link = f"![{alt_text}]({relative_png_path})"
                        self.logger.info(f"PlantUML转换成功: {puml_path} -> {relative_png_path}")
                        return new_link
                    else:
                        self.logger.error(f"PlantUML转换失败，PNG文件不存在: {png_path}")
                else:
                    self.logger.error(f"PlantUML转换失败: {full_puml_path}")
                    
            except Exception as e:
                self.logger.error(f"PlantUML转换异常: {full_puml_path}, 错误: {e}")
            
            # 转换失败时返回原始链接
            return match.group(0)
        
        # 替换所有PlantUML文件链接
        processed_content = plantuml_pattern.sub(replace_plantuml_link, content)
        
        return processed_content, temp_files
            
        # Create a deterministic final output path based on the original input file.
        original_input_path = Path(original_input_file)
        output_path = str(self.output_dir / f"{original_input_path.stem}.docx")
        
        try:
            # 获取模板上下文数据
            context = {
                'project_name': self.project_name or '',
                'title': title,
                'document_no': "P" + datetime.now().strftime("%Y%m%d%H%M%S"),
                'date': datetime.now().strftime("%Y-%m-%d"),
                'author': self.author or '',
                'mobilephone': self.mobilephone or '',
                'email': self.email or ''
            }
            
            self.logger.info(f"使用模板: {template_path}")
            self.logger.info(f"模板上下文: {context}")
            
            # 使用DocxTemplate渲染模板
            doc_tpl = DocxTemplate(template_path)
            doc_tpl.render(context)
            doc_tpl.save(output_path)

            # 加载渲染后的模板文档
            master = Document(output_path)
            
            # 创建composer对象
            composer = Composer(master)
            
            # 加载内容文档
            content_doc = Document(content_path)
            
            # 在模板文档末尾添加连续分节符
            section = master.add_section()
            section.start_type = WD_SECTION_START.CONTINUOUS
            
            # 合并文档，保留样式
            composer.append(content_doc)
            
            # 更新文档属性
            master.core_properties.title = title
            
            # 保存合并后的文档
            composer.save(output_path)
            
            self.logger.info(f"模板处理成功: {output_path}")
            return output_path
            
        except Exception as e:
            self.logger.error(f"模板处理失败: {e}")
            import traceback
            self.logger.error(f"详细错误信息: {traceback.format_exc()}")
            # On failure, return the path to the original content so it can be handled upstream
            return content_path
    
    def _post_process_html(self, html_file: str, processed_md_file: str):
        """后处理HTML文件，添加样式和目录"""
        try:
            # 读取生成的HTML文件
            with open(html_file, 'r', encoding='utf-8') as f:
                html_content = f.read()
            
            # 获取GitHub主题CSS
            theme_css = self._get_github_theme_css()
            
            # 添加锚点ID到标题
            import re
            def add_anchor_to_heading(match):
                level = len(match.group(1))
                title = match.group(2)
                anchor_id = re.sub(r'[^\w\s-]', '', title).strip()
                anchor_id = re.sub(r'[\s_-]+', '-', anchor_id).lower()
                return f'<h{level} id="{anchor_id}">{title}</h{level}>'
            
            html_content = re.sub(r'<h([1-6])>(.*?)</h[1-6]>', add_anchor_to_heading, html_content)
            
            # 插入CSS样式
            css_insert = f'<style>\n{theme_css}\n</style>'
            if '</head>' in html_content:
                html_content = html_content.replace('</head>', f'{css_insert}\n</head>')
            else:
                html_content = f'<head>\n{css_insert}\n</head>\n{html_content}'
            
            # 写回HTML文件
            with open(html_file, 'w', encoding='utf-8') as f:
                f.write(html_content)
                
        except Exception as e:
            self.logger.warning(f"HTML后处理失败: {e}")
    
    def _get_github_theme_css(self) -> str:
        """获取GitHub主题的CSS样式"""
        return """
        body {
            font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif;
            line-height: 1.6;
            color: #24292e;
            max-width: 980px;
            margin: 0 auto;
            padding: 45px;
            background-color: #ffffff;
        }
        h1, h2, h3, h4, h5, h6 {
            margin-top: 24px;
            margin-bottom: 16px;
            font-weight: 600;
            line-height: 1.25;
        }
        h1 { font-size: 2em; border-bottom: 1px solid #eaecef; padding-bottom: 10px; }
        h2 { font-size: 1.5em; border-bottom: 1px solid #eaecef; padding-bottom: 8px; }
        h3 { font-size: 1.25em; }
        h4 { font-size: 1em; }
        h5 { font-size: 0.875em; }
        h6 { font-size: 0.85em; color: #6a737d; }
        p { margin-top: 0; margin-bottom: 16px; }
        blockquote {
            padding: 0 1em;
            color: #6a737d;
            border-left: 0.25em solid #dfe2e5;
            margin: 0 0 16px 0;
        }
        ul, ol { padding-left: 2em; margin-top: 0; margin-bottom: 16px; }
        li { word-wrap: break-all; }
        code {
            padding: 0.2em 0.4em;
            margin: 0;
            font-size: 85%;
            background-color: rgba(27,31,35,0.05);
            border-radius: 3px;
        }
        pre {
            padding: 16px;
            overflow: auto;
            font-size: 85%;
            line-height: 1.45;
            background-color: #f6f8fa;
            border-radius: 3px;
        }
        """
    
    def _optimize_table_column_widths(self, content: str) -> str:
        """
        优化表格列宽分配，特别处理第一列：
        - 如果第一列所有内容显示宽度都小于20（相当于10个汉字），按最宽行设置第一列宽度
        - 否则按平均分配处理
        - 保持原有的 pipe table 格式，不转换为 multiline table
        - 绝对不改变表格内容，只调整格式
        """
        try:
            # 匹配pipe table格式的表格
            table_pattern = re.compile(
                r'(\|[^\n]+\|\n\|[-:\s|]+\|\n(?:\|[^\n]+\|\n?)*)',
                re.MULTILINE
            )
            
            def get_display_width(text):
                """计算文本的显示宽度（中文字符宽度为2，英文字符宽度为1）"""
                width = 0
                for char in text:
                    if '\u4e00' <= char <= '\u9fff':  # 中文字符
                        width += 2
                    elif '\u3000' <= char <= '\u303f':  # 中文标点
                        width += 2
                    elif '\uff00' <= char <= '\uffef':  # 全角字符
                        width += 2
                    else:  # 英文、数字、半角符号等
                        width += 1
                return width
            
            def optimize_table(match):
                table_text = match.group(1).strip()
                lines = table_text.split('\n')
                
                if len(lines) < 3:  # 至少需要标题行、分隔行、数据行
                    return table_text
                
                # 解析表格数据
                header_line = lines[0]
                separator_line = lines[1]
                data_lines = lines[2:]
                
                # 提取列数据
                header_cells = [cell.strip() for cell in header_line.split('|')[1:-1]]
                data_rows = []
                for line in data_lines:
                    if line.strip():
                        cells = [cell.strip() for cell in line.split('|')[1:-1]]
                        if len(cells) == len(header_cells):
                            data_rows.append(cells)
                
                if not header_cells or not data_rows:
                    return table_text
                
                # 检查第一列是否所有内容的显示宽度都小于20（相当于10个汉字）
                first_column_short = True
                first_column_max_width = get_display_width(header_cells[0])
                
                # 检查标题行第一列
                if get_display_width(header_cells[0]) >= 20:
                    first_column_short = False
                
                # 检查数据行第一列
                for row in data_rows:
                    if len(row) > 0:
                        cell_content = row[0]
                        cell_width = get_display_width(cell_content)
                        first_column_max_width = max(first_column_max_width, cell_width)
                        if cell_width >= 20:
                            first_column_short = False
                
                # 计算每列的实际显示宽度
                column_widths = []
                for i in range(len(header_cells)):
                    max_width = get_display_width(header_cells[i])
                    for row in data_rows:
                        if i < len(row):
                            max_width = max(max_width, get_display_width(row[i]))
                    column_widths.append(max_width)
                
                # 智能列宽分配逻辑
                if first_column_short and len(header_cells) > 1:
                    # 第一列内容较短时的特殊处理
                    target_widths = []
                    
                    # 第一列：确保至少8个字符宽度，避免显示问题
                    first_col_width = max(first_column_max_width, 8)
                    target_widths.append(first_col_width)
                    
                    # 其他列：限制最大宽度，避免表格过宽
                    for i in range(1, len(header_cells)):
                        # 限制单列最大宽度为60字符，避免表格过宽
                        col_width = min(column_widths[i], 60)
                        # 确保最小宽度为8字符
                        col_width = max(col_width, 8)
                        target_widths.append(col_width)
                    
                else:
                    # 第一列内容较长时，所有列都限制最大宽度
                    target_widths = []
                    for width in column_widths:
                        # 限制最大宽度为50字符，确保最小宽度为8字符
                        col_width = min(max(width, 8), 50)
                        target_widths.append(col_width)
                
                # 重新构建 pipe table，保持原格式
                result_lines = []
                
                # 标题行
                header_parts = ['|']
                for i, (cell, target_width) in enumerate(zip(header_cells, target_widths)):
                    # 计算需要的空格数来达到目标宽度
                    cell_width = get_display_width(cell)
                    padding = max(1, target_width - cell_width + 2)  # 至少1个空格
                    header_parts.append(f' {cell}' + ' ' * (padding - 1) + '|')
                result_lines.append(''.join(header_parts))
                
                # 分隔行
                separator_parts = ['|']
                for target_width in target_widths:
                    separator_parts.append('-' * (target_width + 2) + '|')
                result_lines.append(''.join(separator_parts))
                
                # 数据行
                for row in data_rows:
                    row_parts = ['|']
                    for i, target_width in enumerate(target_widths):
                        cell = row[i] if i < len(row) else ''
                        cell_width = get_display_width(cell)
                        padding = max(1, target_width - cell_width + 2)  # 至少1个空格
                        row_parts.append(f' {cell}' + ' ' * (padding - 1) + '|')
                    result_lines.append(''.join(row_parts))
                
                return '\n'.join(result_lines)
            
            # 替换所有表格
            optimized_content = table_pattern.sub(optimize_table, content)
            
            # 记录优化的表格数量
            table_count = len(table_pattern.findall(content))
            if table_count > 0:
                self.logger.info(f"优化了 {table_count} 个表格的列宽分配")
            
            return optimized_content
            
        except Exception as e:
            self.logger.warning(f"表格列宽优化失败: {e}")
            return content
    
 