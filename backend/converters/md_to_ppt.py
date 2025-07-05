import os
import re
import glob
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE, MSO_CONNECTOR_TYPE, PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
import pptx.shapes.placeholder
import traceback
from PIL import Image
import shutil
import tempfile
import json
import sys
import time
from typing import List
import logging
from .base_converter import BaseConverter

# 定义目录常量
SCRIPT_DIR = Path(__file__).parent  # tools目录
PROJECT_ROOT = SCRIPT_DIR.parent    # 项目根目录
OUTPUT_DIR = PROJECT_ROOT / 'output'  # 项目根目录下的output目录
TEMP_DIR = Path(tempfile.mkdtemp())  # 临时目录，用于存储生成的图像

class MdToPptConverter(BaseConverter):
    # 默认配置值
    _DEFAULT_TEMPLATE_FILE = str(SCRIPT_DIR / 'templates/template.pptx')

    class PPTTheme:
        """PPT主题定义"""
        # 颜色方案
        COLORS = {
            'primary': RGBColor(0, 114, 198),    # 主色
            'secondary': RGBColor(0, 139, 139),  # 副色
            'accent': RGBColor(255, 140, 0),     # 强调色
            'text': RGBColor(51, 51, 51),        # 正文色
            'background': RGBColor(255, 255, 255) # 背景色
        }
        
        # 字体方案
        FONTS = {
            'title': '微软雅黑',
            'body': '微软雅黑'
        }
        
        # 字号方案
        FONT_SIZES = {
            'title': Pt(40),
            'subtitle': Pt(32),
            'heading1': Pt(28),
            'heading2': Pt(24),
            'body': Pt(18),
            'detail': Pt(14)
        }

    def __init__(self, output_dir: str, **kwargs):
        """
        初始化转换器。
        
        Args:
            output_dir (str): 输出目录。
            **kwargs: 其他配置参数, 如 pptx_template_path, author, project_name 等.
        """
        super().__init__(output_dir, **kwargs)
        
        # 从kwargs获取所有需要的配置
        self.template = self.config.get("pptx_template_path")
        self.use_template_version = self.config.get("use_template_version", True)
        self.use_plain_version = self.config.get("use_plain_version", False)
        self.project_name = self.config.get("project_name")
        self.author = self.config.get("author")
        self.email = self.config.get("email")
        self.mobilephone = self.config.get("mobilephone")

        # 只有在提供了模板路径时才检查它是否存在
        if self.template and not os.path.isabs(self.template):
            possible_template_paths = [
                self.template,
                os.path.join(SCRIPT_DIR, self.template),
                os.path.join(PROJECT_ROOT, self.template)
            ]
            
            found_template = False
            for path in possible_template_paths:
                if os.path.exists(path):
                    self.template = path
                    found_template = True
                    self.logger.info(f"找到模板文件: {self.template}")
                    break
            
            if not found_template:
                self.logger.warning(f"未找到模板文件 {self.template}，尝试过以下路径: {', '.join(possible_template_paths)}")
                # 回退到默认模板
                self.template = self._DEFAULT_TEMPLATE_FILE
                if os.path.exists(self.template):
                    self.logger.info(f"使用默认模板: {self.template}")
                else:
                    self.logger.warning(f"默认模板也不存在: {self.template}，将创建空白演示文稿")
                    self.template = None  # 设置为None，以便后续逻辑可以处理无模板的情况

    def parse_markdown(self, md_file):
        """解析Markdown文件，返回结构化的内容"""
        with open(md_file, 'r', encoding='utf-8') as f:
            content = f.read()
        
        # 提取第一个一级标题作为文档标题
        title = os.path.splitext(os.path.basename(md_file))[0]  # 默认使用文件名作为标题
        first_h1 = re.search(r'^#\s+(.+)$', content, re.MULTILINE)
        if first_h1:
            title = first_h1.group(1).strip()  # 使用第一个一级标题作为标题
        
        # 按标题分割内容
        sections = []
        
        # 提取所有标题行
        heading_pattern = re.compile(r'^(#+)\s+(.+)$', re.MULTILINE)
        headings = heading_pattern.finditer(content)
        
        # 提取标题和内容
        last_pos = 0
        
        # 首先创建一个顶级章节
        sections.append({'level': 1, 'title': title, 'content': []})
        
        for match in headings:
            heading_start = match.start()
            
            # 如果这不是第一个标题，则提取上一个标题到当前标题之间的内容
            if last_pos > 0:
                section_content = content[last_pos:heading_start].strip().split('\n')
                if sections:  # 确保sections不为空
                    sections[-1]['content'] = section_content
            
            # 提取当前标题信息
            level = len(match.group(1))  # '#'的数量表示标题级别
            heading_text = match.group(2).strip()
            
            # 添加新章节
            sections.append({'level': level, 'title': heading_text, 'content': []})
            
            # 更新上一个标题位置
            last_pos = match.end()
        
        # 添加最后一个章节的内容
        if last_pos > 0 and last_pos < len(content):
            section_content = content[last_pos:].strip().split('\n')
            if sections:  # 确保sections不为空
                sections[-1]['content'] = section_content
        
        # 处理没有标题的情况（直接添加内容到第一节）
        if len(sections) == 1 and not sections[0]['content']:
            sections[0]['content'] = content.strip().split('\n')
        
        # 处理文件没有内容的情况
        if not sections:
            sections.append({'level': 1, 'title': title, 'content': []})
        
        return sections

    def save_presentation_with_retry(self, prs, output_file, max_retries=3, delay=1):
        """带重试机制的PPT保存函数"""
        import time
        import errno
        from pathlib import Path
        
        # 确保输出目录存在
        Path(output_file).parent.mkdir(parents=True, exist_ok=True)
        
        for attempt in range(max_retries):
            try:
                # 如果文件存在，先尝试删除
                if os.path.exists(output_file):
                    try:
                        os.remove(output_file)
                    except PermissionError:
                        self.logger.warning(f"Warning: File {output_file} is being used by another process. Retrying...")
                        time.sleep(delay)
                        continue
                
                # 尝试保存文件
                prs.save(output_file)
                self.logger.info(f"Successfully saved {output_file}")
                return True
                
            except PermissionError as e:
                if attempt < max_retries - 1:
                    self.logger.warning(f"Attempt {attempt + 1} failed: File is being used. Retrying in {delay} seconds...")
                    time.sleep(delay)
                else:
                    self.logger.error(f"Error: Could not save file after {max_retries} attempts. File may be in use.")
                    raise
                    
            except Exception as e:
                self.logger.error(f"Unexpected error while saving: {str(e)}")
                raise
        
        return False

    def create_presentation_from_template(self, template_path):
        """根据模板创建演示文稿，如果未提供模板则创建空白演示文稿"""
        if template_path and os.path.exists(template_path):
            self.logger.info(f"INFO: 正在加载模板: {template_path}")
            try:
                prs = Presentation(template_path)
                
                # 清除模板中的所有现有幻灯片，这是一个更稳定的方法
                while len(prs.slides):
                    rId = prs.slides._sldIdLst[0].rId
                    prs.part.drop_rel(rId)
                    del prs.slides._sldIdLst[0]

                self.logger.info(f"INFO: 成功加载并清空模板: {template_path}")
                return prs
            except Exception as e:
                self.logger.error(f"ERROR: 加载模板失败: {template_path}, 错误: {e}")
                self.logger.info("将创建空白演示文稿作为备用方案。")
                return Presentation()
        else:
            self.logger.info("INFO: 未提供模板或模板不存在，正在创建空白演示文稿。")
            return Presentation()

    def find_placeholder_by_type(self, slide, ph_type):
        """在幻灯片中查找特定类型的占位符"""
        for shape in slide.placeholders:
            if shape.placeholder_format.type == ph_type:
                return shape
        return None

    def find_best_layout(self, prs, layout_type='content'):
        """查找最适合的模板布局"""
        # 首先根据名称匹配
        name_patterns = {
            'title': ['标题幻灯片', 'Title Slide', '封面', '首页'],
            'section': ['节标题', 'Section Header', '章节'],
            'content': ['标题和内容', 'Content', 'Title and Content', '内容页']
        }
        
        # 检查模板是否有任何布局
        if len(prs.slide_layouts) == 0:
            self.logger.warning("WARNING: 模板中没有布局，使用空白布局")
            # 创建一个新的演示文稿，获取其默认布局
            temp_prs = Presentation()
            return temp_prs.slide_layouts[0]
        
        # 如果有多个布局，尝试找到匹配的布局
        if len(prs.slide_layouts) > 1:
            # 首先通过名称查找
            patterns = name_patterns.get(layout_type, [])
            for i, layout in enumerate(prs.slide_layouts):
                layout_name = layout.name if hasattr(layout, 'name') else f"布局 {i}"
                for pattern in patterns:
                    if pattern.lower() in layout_name.lower():
                        self.logger.info(f"INFO: 找到匹配的布局 '{layout_name}' 用于 {layout_type} 类型")
                        return layout
            
            # 如果找不到，使用索引位置推断
            # 通常第一个布局是标题页，第二个是内容页
            if layout_type == 'title' and len(prs.slide_layouts) > 0:
                self.logger.info(f"INFO: 使用第一个布局作为标题页: {prs.slide_layouts[0].name if hasattr(prs.slide_layouts[0], 'name') else '未命名布局'}")
                return prs.slide_layouts[0]
            elif layout_type == 'content' and len(prs.slide_layouts) > 1:
                self.logger.info(f"INFO: 使用第二个布局作为内容页: {prs.slide_layouts[1].name if hasattr(prs.slide_layouts[1], 'name') else '未命名布局'}")
                return prs.slide_layouts[1]
        
        # 如果只有一个布局或未找到匹配的布局，使用第一个布局
        self.logger.info(f"INFO: 未找到匹配的 {layout_type} 布局，使用第一个可用布局")
        return prs.slide_layouts[0]

    def extract_logo_blob(self, template_path):
        """从模板第二页提取logo图片的二进制数据"""
        try:
            # 打开模板
            template_prs = Presentation(template_path)
            if len(template_prs.slides) > 1:
                second_slide = template_prs.slides[1]
                
                # 遍历所有形状并记录位置信息
                for shape in second_slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        # 计算位置（转换为英寸）
                        left_in = shape.left/914400
                        top_in = shape.top/914400
                        
                        # 检查是否是目标位置的logo（第4行右侧的图片）
                        if 3.9 <= top_in <= 4.1 and 8.4 <= left_in <= 8.5:
                            self.logger.info(f"Logo found at left={left_in:.2f}in, top={top_in:.2f}in")
                            return shape.image.blob
                            
        except Exception as e:
            self.logger.error(f"Error extracting logo: {str(e)}")
            traceback.print_exc()
        return None

    def create_content_columns(self, prs, slide, content_lines, start_left=0.5, start_top=1.2):
        """创建分栏内容，当列表列数超过三时，每页最多三列分页"""
        # 固定字体大小，确保所有页面保持一致
        TITLE_FONT_SIZE = Pt(24)    # 标题固定24磅
        ITEM_FONT_SIZE = Pt(18)     # 正文固定18磅
        COLUMN_SPACING = 0.3        # 列间距固定0.3英寸
        
        # 定义页面安全区域
        MARGIN_BOTTOM = 0.3  # 底部边距（英寸）
        MARGIN_SIDES = 0.2   # 左右边距（英寸）
        MARGIN_TOP = 0.2     # 顶部边距（英寸）
        
        list_items = []
        current_group = []
        current_group_title = None
        in_code_block = False
        code_block_content = []
        
        for line in content_lines:
            if line.strip():
                # 检查代码块开始或结束标记
                if line.strip().startswith('```'):
                    if in_code_block:
                        if current_group_title and code_block_content:
                            current_group.extend(code_block_content)
                        code_block_content = []
                    in_code_block = not in_code_block
                    continue
                
                if in_code_block:
                    code_block_content.append(line.strip())
                    continue
                
                if re.match(r'^\d+\.\s+\*\*[^*]+\*\*', line) or (line.startswith('- **') and line.endswith('**')):
                    # 保存前一个组，即使它没有内容
                    if current_group_title is not None:
                        list_items.append({'title': current_group_title, 'items': current_group})
                    current_group = []
                    current_group_title = line.strip()
                elif line.startswith('-'):
                    current_group.append(line.strip())
                else:
                    if current_group_title:
                        current_group.append(line.strip())
        
        # 确保最后一个组也被添加，即使它没有内容
        if current_group_title is not None:
            list_items.append({'title': current_group_title, 'items': current_group})
        
        # 如果没有识别到任何列表项但有内容，则使用所有内容作为一个列表项
        if not list_items and content_lines:
            all_items = [line.strip() for line in content_lines if line.strip() and not line.strip().startswith('```')]
            if all_items:
                list_items.append({'title': '', 'items': all_items})
        
        # 分页：每页最多显示3列
        chunks = []
        for i in range(0, len(list_items), 3):
            chunks.append(list_items[i:i+3])
        
        def draw_columns(current_slide, groups):
            # 计算可用宽度（单位：英寸）
            total_width = (prs.slide_width / 914400) - (2 * MARGIN_SIDES)
            available_height = 6 - MARGIN_BOTTOM - MARGIN_TOP
            
            num_columns = len(groups)
            effective_total_width = total_width - (MARGIN_SIDES * 2)
            column_width = (effective_total_width - (COLUMN_SPACING * (num_columns - 1))) / num_columns
            
            for col, group in enumerate(groups):
                # 计算列位置
                left = Inches(MARGIN_SIDES + (column_width + COLUMN_SPACING) * col)
                top = Inches(start_top + MARGIN_TOP)
                width = Inches(column_width)
                height = Inches(available_height)
                
                # 创建文本框
                content_box = current_slide.shapes.add_textbox(left, top, width, height)
                content_box.line.width = 0
                tf = content_box.text_frame
                tf.word_wrap = True
                tf.auto_size = MSO_AUTO_SIZE.NONE
                tf.margin_bottom = 0
                tf.margin_left = Pt(2)
                tf.margin_right = Pt(2)
                tf.margin_top = 0
                
                # 添加标题 - 即使没有项目内容，也显示标题
                if group['title']:
                    p = tf.add_paragraph()
                    title_text = re.sub(r'\*\*', '', group['title'])
                    title_text = re.sub(r'^-\s*', '', title_text)
                    title_text = re.sub(r'^\d+\.\s*', '', title_text)
                    p.text = title_text
                    p.font.size = TITLE_FONT_SIZE
                    p.font.bold = True
                    p.space_after = Pt(12)
                    p.alignment = PP_ALIGN.LEFT
                
                # 添加列表项
                for item in group['items']:
                    p = tf.add_paragraph()
                    item_text = item
                    if item.startswith('-'):
                        item_text = re.sub(r'^-\s*', '', item_text)
                    p.text = item_text
                    p.font.size = ITEM_FONT_SIZE
                    p.space_before = Pt(6)
                    p.alignment = PP_ALIGN.LEFT
                    
                    if '**' in item_text:
                        p.text = re.sub(r'\*\*', '', p.text)
                        p.font.bold = True
                
                # 确保至少有一行文本，即使没有内容
                if group['title'] and not group['items']:
                    p = tf.add_paragraph()
                    p.text = " "  # 添加一个空格作为占位符
                    p.font.size = ITEM_FONT_SIZE
                    p.space_before = Pt(6)
                    p.alignment = PP_ALIGN.LEFT
                
                p.font.name = "微软雅黑"
        
        # 处理第一页
        draw_columns(slide, chunks[0])
        
        # 处理后续分页
        for chunk in chunks[1:]:
            # 创建新幻灯片
            new_slide = prs.slides.add_slide(prs.slide_layouts[2])
            
            # 复制标题和装饰元素
            # 添加标题左侧的红色方块
            square_size = Inches(0.3)
            red_square = new_slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 
                Inches(0.5),
                Inches(0.5),
                square_size, 
                square_size
            )
            red_square.fill.solid()
            red_square.fill.fore_color.rgb = RGBColor(192, 0, 0)
            red_square.line.fill.background()
            
            # 添加红色横线
            line = new_slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0),
                Inches(1.2),
                prs.slide_width,
                Inches(0.03)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor(192, 0, 0)
            line.line.fill.background()
            
            # 绘制内容
            draw_columns(new_slide, chunk)
        
        return

    def process_md_file(self, md_file, output_file, use_template=True):
        """处理单个Markdown文件并输出到指定路径（源于原始脚本的核心逻辑）"""
        template_path = self.template if use_template else None
        if use_template and not os.path.exists(str(template_path)):
            self.logger.warning(f"警告: 模板文件不存在 {template_path}")
            
        try:
            # 创建演示文稿
            prs = self.create_presentation_from_template(template_path)
            
            # 解析并转换内容
            sections = self.parse_markdown(md_file)
            if not sections:
                self.logger.warning(f"警告: 在 {md_file} 中未找到内容")
                raise ValueError("Markdown文件中未找到内容")
            
            # 获取 Markdown 文件所在目录
            md_dir = os.path.dirname(os.path.abspath(md_file))
            
            # --- 1. 创建标题页 ---
            title_layout = self.find_best_layout(prs, 'title')
            slide = prs.slides.add_slide(title_layout)
            
            # 添加标题文本框
            left = Inches(1); top = Inches(2); width = Inches(14); height = Inches(2)
            title_box = slide.shapes.add_textbox(left, top, width, height)
            title_frame = title_box.text_frame
            title_frame.text = self.project_name or sections[0]['title']
            
            # 设置标题格式
            p = title_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT; p.font.size = Pt(44); p.font.bold = True
            
            # 提取logo
            logo_blob = None
            if use_template and os.path.exists(str(template_path)):
                logo_blob = self.extract_logo_blob(template_path)
            
            # --- 2. 循环创建内容页 ---
            content_layout = self.find_best_layout(prs, 'content')
            section_layout = self.find_best_layout(prs, 'section') 
            if not section_layout: section_layout = title_layout

            for section in sections:
                if section['level'] == 1:
                    continue
                
                is_lower_level_title = section['level'] >= 3
                has_decimal_number = bool(re.match(r'^\d+\.\d+', section['title'].strip()))
                
                if section['level'] == 2 and not (is_lower_level_title or has_decimal_number):
                    # 二级标题处理 - 创建节标题页
                    slide = prs.slides.add_slide(section_layout)
                    left = Inches(1); top = Inches(2); width = Inches(14); height = Inches(2)
                    title_box = slide.shapes.add_textbox(left, top, width, height)
                    title_frame = title_box.text_frame
                    title_frame.text = section['title']
                    p = title_frame.paragraphs[0]
                    p.alignment = PP_ALIGN.LEFT; p.font.size = Pt(40); p.font.bold = True
                    
                    if section['content'] and any(line.strip() for line in section['content']):
                        slide = prs.slides.add_slide(content_layout)
                        self._add_slide_background_and_decorations(prs, slide, logo_blob, title=False)
                        self._populate_slide_content(prs, slide, section, md_dir, start_top=0.9)

                elif section['content']:
                    slide = prs.slides.add_slide(content_layout)
                    self._add_slide_background_and_decorations(prs, slide, logo_blob, title=True, section_title=section['title'])
                    self._populate_slide_content(prs, slide, section, md_dir, start_top=1.4)

            self.save_presentation_with_retry(prs, output_file)
            self.logger.info(f"成功转换 {md_file} 为 {output_file}")
            return True
            
        except Exception as e:
            self.logger.error(f"处理 {md_file} 时发生严重错误: {e}", exc_info=True)
            return False

    def _add_slide_background_and_decorations(self, prs, slide, logo_blob, title=True, section_title=""):
        """辅助函数：为幻灯片添加背景、logo和装饰线条"""
        # 创建纯白色背景
        left = Inches(0); top = Inches(0); width = prs.slide_width; height = prs.slide_height
        try:
            bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
            bg_shape.fill.solid()
            bg_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
            bg_shape.line.fill.background()
            # 尝试将形状发送到底层
            slide.shapes._spTree.insert(2, bg_shape._element)
        except Exception as e:
            self.logger.warning(f"设置背景时出错: {e}")

        # 添加logo到右上角
        if logo_blob:
            from io import BytesIO
            logo_stream = BytesIO(logo_blob)
            logo_height = Inches(0.4)
            try:
                with Image.open(logo_stream) as img:
                    aspect_ratio = img.width / img.height
                    logo_width = logo_height * aspect_ratio
                    logo_left = prs.slide_width - logo_width - Inches(0.3)
                    logo_top = Inches(0.3)
                    slide.shapes.add_picture(BytesIO(logo_blob), logo_left, logo_top, logo_width, logo_height)
            except Exception as e:
                 self.logger.warning(f"添加Logo时出错: {e}")

        # 添加底部红线
        left = Inches(0); top = prs.slide_height - Inches(0.2); width = prs.slide_width; height = Inches(0.2)
        bottom_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        bottom_line.fill.solid()
        bottom_line.fill.fore_color.rgb = RGBColor(192, 0, 0)
        bottom_line.line.fill.background()

        # 添加标题和装饰
        line_top = Inches(1.2)
        if title:
            left = Inches(1.3); top = Inches(0.5); width = Inches(8); height = Inches(0.6)
            title_box = slide.shapes.add_textbox(left, top, width, height)
            p = title_box.text_frame.add_paragraph()
            p.text = section_title
            p.alignment = PP_ALIGN.LEFT; p.font.size = Pt(24); p.font.bold = True
            
            left = Inches(0.8); top = Inches(0.5); width = Inches(0.3); height = Inches(0.3)
            rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
            rect.fill.solid(); rect.fill.fore_color.rgb = RGBColor(192, 0, 0); rect.line.fill.background()
        else:
            line_top = Inches(0.7)

        left = Inches(0); top = line_top; width = prs.slide_width; height = Inches(0.03)
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(192, 0, 0)
        line.line.fill.background()

    def _populate_slide_content(self, prs, slide, section, md_dir, start_top):
        """辅助函数：填充幻灯片内容，包括文本、列表和图片"""
        content_lines = []
        for line in section['content']:
            img_match = re.search(r'!\[[^\]]*\]\((.*?)\)', line)
            if img_match:
                img_path = img_match.group(1)
                if not os.path.isabs(img_path): img_path = os.path.join(md_dir, img_path)
                if os.path.exists(img_path): content_lines.append(f"<image:{img_path}>")
                else: self.logger.warning(f"图片未找到: {img_path}")
            elif line.strip(): content_lines.append(line.strip())

        is_list = any(line.strip().startswith('-') or re.match(r'^\d+\.', line.strip()) for line in content_lines if not line.startswith("<image:"))

        text_lines = [line for line in content_lines if not line.startswith("<image:")]
        image_lines = [line for line in content_lines if line.startswith("<image:")]

        img_y_pos = start_top

        if is_list:
            self.create_content_columns(prs, slide, text_lines, start_left=0.8, start_top=start_top)
            num_rows = max(len(text_lines) / 3, 1) if text_lines else 1
            img_y_pos = start_top + Inches(num_rows * 0.5) + Inches(0.5)
        else:
            content_left = Inches(0.8); content_top = Inches(start_top)
            content_width = prs.slide_width - Inches(1.6); content_height = prs.slide_height - Inches(start_top + 0.5)
            if text_lines:
                text_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
                text_frame = text_box.text_frame; text_frame.word_wrap = True
                for line in text_lines:
                    p = text_frame.add_paragraph()
                    p.text = line; p.font.size = Pt(18); p.alignment = PP_ALIGN.LEFT; p.font.name = "微软雅黑"
                
                current_text_height = Inches(len(text_frame.paragraphs) * 0.3)
                img_y_pos = content_top + current_text_height + Inches(0.2)

        for line in image_lines:
            img_path = line.split(":", 1)[1]
            try:
                with Image.open(img_path) as img:
                    aspect_ratio = img.width / img.height
                    img_max_width = prs.slide_width - Inches(2)
                    img_max_height = Inches(3.0)
                    img_width = min(img_max_width, img_max_height * aspect_ratio)
                    img_height = min(img_max_height, img_width / aspect_ratio)
                    
                    if img_y_pos + img_height > prs.slide_height - Inches(0.5):
                         img_y_pos = start_top

                    slide.shapes.add_picture(img_path, Inches(1), img_y_pos, width=img_width, height=img_height)
                    img_y_pos += img_height + Inches(0.2)
            except Exception as e:
                self.logger.warning(f"添加图片时出错 {img_path}: {e}")

    def convert(self, input_path: str) -> List[str]:
        """
        将Markdown文件或目录转换为PPT。
        此方法实现了BaseConverter的抽象接口，并整合了所有高级转换逻辑。
        它能够处理单个文件或目录，并根据配置生成带模板和不带模板的多个版本。

        Args:
            input_path: 输入文件或目录的路径。

        Returns:
            List[str]: 生成的PPT文件路径列表。
        """
        supported_extensions = ['.md', '.markdown']
        if not self._is_valid_input(input_path, supported_extensions):
            self.logger.error(f"无效的输入路径或文件类型: {input_path}")
            return []

        files_to_convert = []
        if os.path.isfile(input_path):
            files_to_convert.append(input_path)
        else:  # it's a directory
            files_to_convert.extend(self._get_files_by_extension(input_path, supported_extensions))
        
        if not files_to_convert:
            self.logger.warning(f"在 {input_path} 中未找到要转换的Markdown文件。")
            return []

        generated_files = []
        for md_file in files_to_convert:
            # 1. 生成不使用模板的版本
            if self.use_plain_version:
                output_file_plain = self._generate_output_path(md_file, ".pptx")
                try:
                    self.logger.info(f"正在生成普通版本: {output_file_plain}")
                    if self.process_md_file(md_file, str(output_file_plain), use_template=False):
                        generated_files.append(str(output_file_plain))
                except Exception as e:
                    self.logger.error(f"转换 {md_file} 为普通版本时出错: {str(e)}", exc_info=True)

            # 2. 生成使用模板的版本
            if self.use_template_version:
                # 当同时生成两个版本时，为模板版本添加后缀以区分
                if self.use_plain_version:
                    output_file_template = self._generate_output_path(md_file, ".template.pptx")
                else:
                    output_file_template = self._generate_output_path(md_file, ".pptx")
                
                try:
                    self.logger.info(f"正在生成模板版本: {output_file_template}")
                    if self.process_md_file(md_file, str(output_file_template), use_template=True):
                        generated_files.append(str(output_file_template))
                except Exception as e:
                    self.logger.error(f"转换 {md_file} 为模板版本时出错: {str(e)}", exc_info=True)
        
        return generated_files

    def convert_md_to_ppt(self, input_file: str, output_dir: str = None) -> str:
        """
        将Markdown文件转换为PPT
        
        Args:
            input_file (str): 输入的Markdown文件路径
            output_dir (str, optional): 输出目录
            
        Returns:
            str: 输出文件路径，成功时返回路径，失败时返回None
        """
        # 如果指定了输出目录，使用它；否则使用项目根目录下的output目录
        output_path = Path(output_dir) if output_dir else OUTPUT_DIR
        
        # 确保输出目录存在
        output_path.mkdir(parents=True, exist_ok=True)
        
        # 创建输出文件名
        input_path = Path(input_file)
        output_file = str(output_path / f"{input_path.stem}.pptx")
        
        try:
            # 使用全局模板路径
            template_to_use = self.template
                
            # 调用process_md_file函数
            if self.process_md_file(input_file, output_file, use_template=True):
                return output_file
            else:
                return None
        
        except Exception as e:
            self.logger.error(f"转换 {input_file} 时出错: {str(e)}")
            traceback.print_exc()
            return None
